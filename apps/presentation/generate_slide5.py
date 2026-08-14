import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_BG_DARK = RGBColor(10, 11, 12)     # Very Dark Charcoal/Black (#0A0B0C)
COLOR_CARD = RGBColor(22, 24, 26)        # Dark Card Background (#16181A)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_WHITE = RGBColor(245, 245, 247)    # Off-white text (#F5F5F7)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(120, 120, 125)    # Muted Gray-green (#78787D)
COLOR_BORDER_THIN = RGBColor(220, 222, 218) # Thin card borders (#DCdeda)
COLOR_FOREST = RGBColor(31, 77, 58)        # Deep Forest Green (#1F4D3A)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_connector(slide, x1, y1, x2, y2, color, width_pt=1.0):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn

# Helper vector shapes for benefit icons
def draw_icon_network(slide, cx, cy):
    # Center node
    draw_oval(slide, cx, cy, 0.02, COLOR_ORANGE)
    # 5 surrounding nodes connected to center
    r_node = 0.11
    for i in range(5):
        angle = math.radians(i * 72)
        nx = cx + r_node * math.cos(angle)
        ny = cy + r_node * math.sin(angle)
        draw_connector(slide, cx, cy, nx, ny, COLOR_CHARCOAL, 0.8)
        draw_oval(slide, nx, ny, 0.022, COLOR_CHARCOAL)

def draw_icon_chart(slide, cx, cy):
    # 3 chart pillars
    draw_line(slide, cx - 0.08, cy + 0.04, cx - 0.04, cy + 0.10, COLOR_CHARCOAL)
    draw_line(slide, cx - 0.02, cy + 0.00, cx + 0.02, cy + 0.10, COLOR_CHARCOAL)
    draw_line(slide, cx + 0.04, cy - 0.04, cx + 0.08, cy + 0.10, COLOR_CHARCOAL)
    # Rising arrow in orange
    draw_connector(slide, cx - 0.09, cy + 0.06, cx + 0.09, cy - 0.08, COLOR_ORANGE, 1.2)
    # Arrow tip
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx + 0.05), Inches(cy - 0.11), Inches(0.06), Inches(0.06)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_ORANGE
    arrow.line.fill.background()

def draw_icon_security(slide, cx, cy):
    # Shield outline
    w = 0.06
    h = 0.07
    builder = slide.shapes.build_freeform(Inches(cx - w), Inches(cy - h))
    vertices = [
        (Inches(cx + w), Inches(cy - h)),
        (Inches(cx + w), Inches(cy)),
        (Inches(cx), Inches(cy + h)),
        (Inches(cx - w), Inches(cy)),
        (Inches(cx - w), Inches(cy - h))
    ]
    builder.add_line_segments(vertices, close=True)
    shield = builder.convert_to_shape()
    shield.fill.background()
    shield.line.color.rgb = COLOR_CHARCOAL
    shield.line.width = Pt(1)
    
    # Tiny orange lock inside
    draw_oval(slide, cx, cy - 0.01, 0.02, COLOR_ORANGE)
    body = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.03), Inches(cy + 0.01), Inches(0.06), Inches(0.05)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = COLOR_ORANGE
    body.line.fill.background()

def draw_icon_efficiency(slide, cx, cy):
    # Stacked cylinders on the left
    for dy in [-0.05, 0, 0.05]:
        c = draw_oval(slide, cx - 0.04, cy + dy, 0.05, COLOR_BG_LIGHT, line_color=COLOR_CHARCOAL, line_width_pt=1)
        c.fill.background()
        
    # Orange dollar circle on bottom right
    draw_oval(slide, cx + 0.06, cy + 0.05, 0.04, COLOR_ORANGE)
    # Simple "$" inside
    dl_box = slide.shapes.add_textbox(Inches(cx + 0.03), Inches(cy + 0.015), Inches(0.06), Inches(0.07))
    tf = dl_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "$"
    p.font.name = 'Inter'
    p.font.size = Pt(6)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

def draw_icon_scale(slide, cx, cy):
    # Dartboard concentric ovals
    draw_oval(slide, cx, cy, 0.11, COLOR_BG_LIGHT, line_color=COLOR_CHARCOAL, line_width_pt=1).fill.background()
    draw_oval(slide, cx, cy, 0.05, COLOR_BG_LIGHT, line_color=COLOR_CHARCOAL, line_width_pt=1).fill.background()
    
    # Arrow in orange
    draw_connector(slide, cx + 0.09, cy - 0.09, cx, cy, COLOR_ORANGE, 1.2)
    # Arrow fletchings/tail
    draw_line(slide, cx + 0.08, cy - 0.10, cx + 0.11, cy - 0.07, COLOR_ORANGE, 1)

def draw_icon_ecosystem(slide, cx, cy):
    # Planet sphere
    draw_oval(slide, cx, cy, 0.09, COLOR_BG_LIGHT, line_color=COLOR_CHARCOAL, line_width_pt=1).fill.background()
    # Orbits/Lines
    draw_connector(slide, cx - 0.11, cy + 0.02, cx + 0.11, cy - 0.02, COLOR_CHARCOAL, 0.8)
    draw_connector(slide, cx - 0.02, cy - 0.11, cx + 0.02, cy + 0.11, COLOR_CHARCOAL, 0.8)
    # Nodes on orbits
    draw_oval(slide, cx - 0.06, cy + 0.04, 0.02, COLOR_ORANGE)
    draw_oval(slide, cx + 0.06, cy - 0.04, 0.02, COLOR_ORANGE)

def create_slide_5():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Slide Index "05"
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "05"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(22)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Vertical divider on top left
    draw_line(slide, 1.45, 0.42, 1.46, 0.78, COLOR_BORDER_THIN)
    
    # Top left subheader
    sh_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.4), Inches(3.0), Inches(0.5))
    tf_sh = sh_box.text_frame
    tf_sh.word_wrap = True
    p_sh = tf_sh.paragraphs[0]
    p_sh.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_sh.font.name = 'Inter'
    p_sh.font.size = Pt(8.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_MUTED
    p_sh.line_spacing = 1.15
    
    # 2. Right Top Section Header: KEY BENEFITS
    draw_line(slide, 10.3, 0.42, 10.315, 0.78, COLOR_ORANGE)
    sol_hdr_box = slide.shapes.add_textbox(Inches(10.45), Inches(0.42), Inches(2.0), Inches(0.4))
    p_shdr = sol_hdr_box.text_frame.paragraphs[0]
    p_shdr.text = "KEY BENEFITS"
    p_shdr.font.name = 'Inter'
    p_shdr.font.size = Pt(10)
    p_shdr.font.bold = True
    p_shdr.font.color.rgb = COLOR_CHARCOAL
    
    # 3. Left Title Box (Compact vertical dimensions)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4.2), Inches(3.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    
    p_lbl = tf_title.paragraphs[0]
    p_lbl.text = "KEY BENEFITS"
    p_lbl.font.name = 'Inter'
    p_lbl.font.size = Pt(10.5)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = COLOR_ORANGE
    p_lbl.space_after = Pt(10)
    
    # Sora heading runs (Set size to 32pt and strict line height to Pt(36) for compact styling)
    p_h = tf_title.add_paragraph()
    p_h.line_spacing = Pt(36)
    p_h.space_after = Pt(12)
    
    r_u = p_h.add_run()
    r_u.text = "UNLOCKING THE\nPOWER OF\n"
    r_u.font.name = 'Sora'
    r_u.font.size = Pt(32)
    r_u.font.bold = True
    r_u.font.color.rgb = COLOR_CHARCOAL
    
    r_c = p_h.add_run()
    r_c.text = "COLLABORATIVE\nINTELLIGENCE"
    r_c.font.name = 'Sora'
    r_c.font.size = Pt(32)
    r_c.font.bold = True
    r_c.font.color.rgb = COLOR_ORANGE
    
    # Horizontal line under title
    draw_line(slide, 0.8, 5.4, 1.3, 5.42, COLOR_ORANGE)
    
    # Description paragraph
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.65), Inches(3.6), Inches(1.1))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = 0
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "Our solution transforms isolated AI systems into a connected ecosystem that drives measurable impact across organizations and industries."
    p_desc.font.name = 'Inter'
    p_desc.font.size = Pt(12.5)
    p_desc.font.color.rgb = COLOR_CHARCOAL
    p_desc.line_spacing = 1.3
    
    # 4. Right Side 2x3 Cards Grid (Optimized spacings and margins)
    cards_data = [
        # (col, row, title, desc, icon_type)
        (0, 0, "STRONGER\nCOLLABORATION", "Break down silos and enable seamless collaboration across organizations.", 'network'),
        (1, 0, "ACCELERATED\nINNOVATION", "Combine diverse data, models, and expertise to drive faster breakthroughs.", 'chart'),
        (2, 0, "ENHANCED\nSECURITY & TRUST", "Built with privacy-preserving technology and zero-trust architecture.", 'security'),
        (0, 1, "GREATER\nEFFICIENCY", "Reduce duplication, optimize resources, and lower operational costs.", 'efficiency'),
        (1, 1, "IMPACT AT\nSCALE", "Solve bigger challenges together and create lasting value for society.", 'scale'),
        (2, 1, "FUTURE-READY\nECOSYSTEM", "A scalable foundation designed to evolve with emerging technologies.", 'ecosystem')
    ]
    
    for col, row, title, desc, icon_type in cards_data:
        left = 5.2 + col * 2.6
        top = 1.1 + row * 2.85
        width = 2.4
        height = 2.6
        
        # Base Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_BG_LIGHT
        card.line.color.rgb = COLOR_BORDER_THIN
        card.line.width = Pt(1)
        
        # Icon outer circle
        cx, cy = left + 1.2, top + 0.55
        draw_oval(slide, cx, cy, 0.325, COLOR_BG_LIGHT, line_color=COLOR_BORDER_THIN, line_width_pt=0.8).fill.background()
        
        # Render internal icon
        if icon_type == 'network':
            draw_icon_network(slide, cx, cy)
        elif icon_type == 'chart':
            draw_icon_chart(slide, cx, cy)
        elif icon_type == 'security':
            draw_icon_security(slide, cx, cy)
        elif icon_type == 'efficiency':
            draw_icon_efficiency(slide, cx, cy)
        elif icon_type == 'scale':
            draw_icon_scale(slide, cx, cy)
        elif icon_type == 'ecosystem':
            draw_icon_ecosystem(slide, cx, cy)
            
        # Card title (With precise margins and space parameters)
        t_box = slide.shapes.add_textbox(Inches(left + 0.05), Inches(top + 1.05), Inches(2.3), Inches(0.55))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Inter'
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_CHARCOAL
        p_t.alignment = PP_ALIGN.CENTER
        p_t.line_spacing = Pt(12.5)
        p_t.space_after = Pt(4)
        
        # Card description (With compact point-based line height to prevent word overlap)
        d_box = slide.shapes.add_textbox(Inches(left + 0.05), Inches(top + 1.68), Inches(2.3), Inches(0.8))
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_right = tf_d.margin_top = tf_d.margin_bottom = 0
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = COLOR_MUTED
        p_d.alignment = PP_ALIGN.CENTER
        p_d.line_spacing = Pt(11.5)

    # 5. Bottom Horizontal Black Band
    quote_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0), Inches(7.0), Inches(13.333), Inches(0.5)
    )
    quote_band.fill.solid()
    quote_band.fill.fore_color.rgb = COLOR_BG_DARK
    quote_band.line.fill.background()
    
    # Left highlight
    draw_oval(slide, 0.8, 7.25, 0.16, COLOR_BG_DARK, line_color=COLOR_ORANGE, line_width_pt=1)
    draw_oval(slide, 0.8, 7.25, 0.08, COLOR_ORANGE)
    draw_line(slide, 1.2, 7.1, 1.215, 7.4, COLOR_ORANGE)
    
    quote_box = slide.shapes.add_textbox(Inches(1.35), Inches(7.1), Inches(7.5), Inches(0.35))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = 0
    p_q = tf_q.paragraphs[0]
    p_q.text = ""
    r1 = p_q.add_run()
    r1.text = "Together, we can turn independent intelligence into "
    r1.font.name = 'Inter'
    r1.font.size = Pt(10.5)
    r1.font.color.rgb = COLOR_WHITE
    r2 = p_q.add_run()
    r2.text = "collective impact"
    r2.font.name = 'Inter'
    r2.font.size = Pt(10.5)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_ORANGE
    r3 = p_q.add_run()
    r3.text = "."
    r3.font.name = 'Inter'
    r3.font.size = Pt(10.5)
    r3.font.color.rgb = COLOR_WHITE
    
    # Right sub-label
    draw_line(slide, 9.8, 7.1, 9.815, 7.4, COLOR_ORANGE)
    right_lbl_box = slide.shapes.add_textbox(Inches(10.0), Inches(7.08), Inches(3.0), Inches(0.4))
    tf_rl = right_lbl_box.text_frame
    tf_rl.word_wrap = True
    tf_rl.margin_left = tf_rl.margin_right = tf_rl.margin_top = tf_rl.margin_bottom = 0
    p_rl = tf_rl.paragraphs[0]
    p_rl.text = "CONNECTED INTELLIGENCE.\nSTRONGER TOGETHER."
    p_rl.font.name = 'Inter'
    p_rl.font.size = Pt(8.5)
    p_rl.font.bold = True
    p_rl.font.color.rgb = COLOR_WHITE
    p_rl.line_spacing = 1.05

    # Save
    output_filename = "presentation_slide5_final.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    try:
        prs.save(output_path)
        print(f"Slide 5 presentation successfully generated at: {output_path}")
    except PermissionError:
        fallback_path = os.path.join(output_dir, "presentation_slide5_fixed.pptx")
        prs.save(fallback_path)
        print(f"[!] WARNING: '{output_path}' was locked/open. Saved fallback to: {fallback_path}")

if __name__ == "__main__":
    create_slide_5()
