import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BG = RGBColor(10, 11, 12)          # Very Dark Charcoal/Black (#0A0B0C)
COLOR_CARD = RGBColor(22, 24, 26)        # Dark Card Background (#16181A)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_WHITE = RGBColor(245, 245, 247)    # Off-white text (#F5F5F7)
COLOR_GRAY = RGBColor(160, 160, 165)     # Muted gray text (#A0A0A5)
COLOR_DIVIDER = RGBColor(36, 38, 40)     # Dark Gray Divider (#242628)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1), Inches(y2 - y1 if (y2 - y1) > 0 else 0.01)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_user_icon(slide, cx, cy):
    # Head
    draw_oval(slide, cx, cy - 0.06, 0.045, COLOR_ORANGE)
    # Torso (rounded rectangle)
    torso = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(cx - 0.09), Inches(cy + 0.03), Inches(0.18), Inches(0.09)
    )
    torso.fill.solid()
    torso.fill.fore_color.rgb = COLOR_ORANGE
    torso.line.fill.background()

def draw_guide_icon(slide, cx, cy):
    # Outer credential card box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(cx - 0.09), Inches(cy - 0.09), Inches(0.18), Inches(0.18)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CARD
    box.line.color.rgb = COLOR_ORANGE
    box.line.width = Pt(1)
    
    # Inner head-and-torso details
    draw_oval(slide, cx, cy - 0.03, 0.03, COLOR_ORANGE)
    
    torso = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(cx - 0.06), Inches(cy + 0.02), Inches(0.12), Inches(0.05)
    )
    torso.fill.solid()
    torso.fill.fore_color.rgb = COLOR_ORANGE
    torso.line.fill.background()

def draw_dept_icon(slide, cx, cy):
    # Columns / pillars (3 lines)
    draw_line(slide, cx - 0.07, cy - 0.03, cx - 0.05, cy + 0.05, COLOR_ORANGE, 1)
    draw_line(slide, cx - 0.01, cy - 0.03, cx + 0.01, cy + 0.05, COLOR_ORANGE, 1)
    draw_line(slide, cx + 0.05, cy - 0.03, cx + 0.07, cy + 0.05, COLOR_ORANGE, 1)
    
    # Base
    base = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(cx - 0.1), Inches(cy + 0.05), Inches(0.2), Inches(0.02)
    )
    base.fill.solid()
    base.fill.fore_color.rgb = COLOR_ORANGE
    base.line.fill.background()
    
    # Roof (Triangle)
    roof = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, 
        Inches(cx - 0.11), Inches(cy - 0.1), Inches(0.22), Inches(0.07)
    )
    roof.fill.solid()
    roof.fill.fore_color.rgb = COLOR_ORANGE
    roof.line.fill.background()

def create_slide_1():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Left border accent line
    draw_line(slide, 0.35, 0.0, 0.38, 7.5, COLOR_ORANGE)
    
    # 2. Slide Number "01"
    num_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    tf_num.word_wrap = True
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    p_num = tf_num.paragraphs[0]
    p_num.text = "01"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(18)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # 3. Main Title Box
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(7.1), Inches(3.9))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    
    p_title = tf_title.paragraphs[0]
    p_title.line_spacing = 0.95
    p_title.space_after = Pt(12)
    
    # Run 1: Connecting Intelligence
    r1 = p_title.add_run()
    r1.text = "Connecting\nIntelligence\n"
    r1.font.name = 'Sora'
    r1.font.size = Pt(54)
    r1.font.bold = True
    r1.font.color.rgb = COLOR_WHITE
    
    # Run 2: Across Organizations
    r2 = p_title.add_run()
    r2.text = "Across\nOrganizations"
    r2.font.name = 'Sora'
    r2.font.size = Pt(54)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_ORANGE
    
    # 4. Horizontal Divider Line under Title
    draw_line(slide, 0.7, 5.25, 2.7, 5.28, COLOR_ORANGE)
    
    # 5. Subtitle Text Box
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.65), Inches(7.1), Inches(1.1))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "A DISTRIBUTED INTELLIGENCE INFRASTRUCTURE\nFOR SECURE AI COLLABORATION"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = COLOR_WHITE
    p_sub.font.bold = True
    p_sub.line_spacing = 1.3
    
    # 6. Right Side Rounded Rectangle Card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(8.2), Inches(1.0), Inches(4.3), Inches(5.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.fill.background() # No border line
    
    # SECTION 1: PRESENTED BY
    # Icon circle
    draw_oval(slide, 8.82, 1.6, 0.22, COLOR_CARD, line_color=COLOR_ORANGE, line_width_pt=1.2)
    draw_user_icon(slide, 8.82, 1.6)
    
    # Header Label
    lbl1_box = slide.shapes.add_textbox(Inches(9.2), Inches(1.45), Inches(3.0), Inches(0.3))
    tf_lbl1 = lbl1_box.text_frame
    p_lbl1 = tf_lbl1.paragraphs[0]
    p_lbl1.text = "PRESENTED BY"
    p_lbl1.font.name = 'Inter'
    p_lbl1.font.size = Pt(9.5)
    p_lbl1.font.bold = True
    p_lbl1.font.color.rgb = COLOR_ORANGE
    
    # Names list
    names_box = slide.shapes.add_textbox(Inches(8.6), Inches(2.0), Inches(3.6), Inches(1.1))
    tf_names = names_box.text_frame
    tf_names.word_wrap = True
    tf_names.margin_left = tf_names.margin_right = tf_names.margin_top = tf_names.margin_bottom = 0
    p_names = tf_names.paragraphs[0]
    p_names.text = "Gokul S\nGautham Krishna R Nair\nSreenidhi V\nT Maheswaran"
    p_names.font.name = 'Inter'
    p_names.font.size = Pt(12)
    p_names.font.color.rgb = COLOR_WHITE
    p_names.line_spacing = 1.3
    
    # Divider line 1
    draw_line(slide, 8.6, 3.3, 12.1, 3.31, COLOR_DIVIDER)
    
    # SECTION 2: GUIDE
    # Icon circle
    draw_oval(slide, 8.82, 3.8, 0.22, COLOR_CARD, line_color=COLOR_ORANGE, line_width_pt=1.2)
    draw_guide_icon(slide, 8.82, 3.8)
    
    # Header Label
    lbl2_box = slide.shapes.add_textbox(Inches(9.2), Inches(3.65), Inches(3.0), Inches(0.3))
    tf_lbl2 = lbl2_box.text_frame
    p_lbl2 = tf_lbl2.paragraphs[0]
    p_lbl2.text = "GUIDE"
    p_lbl2.font.name = 'Inter'
    p_lbl2.font.size = Pt(9.5)
    p_lbl2.font.bold = True
    p_lbl2.font.color.rgb = COLOR_ORANGE
    
    # Guide Name
    guide_box = slide.shapes.add_textbox(Inches(8.6), Inches(4.2), Inches(3.6), Inches(0.4))
    tf_guide = guide_box.text_frame
    tf_guide.word_wrap = True
    tf_guide.margin_left = tf_guide.margin_right = tf_guide.margin_top = tf_guide.margin_bottom = 0
    p_guide = tf_guide.paragraphs[0]
    p_guide.text = "Assistant Professor Nisha Soman"
    p_guide.font.name = 'Inter'
    p_guide.font.size = Pt(12)
    p_guide.font.color.rgb = COLOR_WHITE
    
    # Divider line 2
    draw_line(slide, 8.6, 4.8, 12.1, 4.81, COLOR_DIVIDER)
    
    # SECTION 3: DEPARTMENT OF ARTIFICIAL INTELLIGENCE...
    # Icon circle
    draw_oval(slide, 8.82, 5.2, 0.22, COLOR_CARD, line_color=COLOR_ORANGE, line_width_pt=1.2)
    draw_dept_icon(slide, 8.82, 5.2)
    
    # Header Label
    lbl3_box = slide.shapes.add_textbox(Inches(9.2), Inches(4.95), Inches(3.0), Inches(0.6))
    tf_lbl3 = lbl3_box.text_frame
    tf_lbl3.word_wrap = True
    p_lbl3 = tf_lbl3.paragraphs[0]
    p_lbl3.text = "DEPARTMENT OF ARTIFICIAL\nINTELLIGENCE AND MACHINE LEARNING"
    p_lbl3.font.name = 'Inter'
    p_lbl3.font.size = Pt(8.5)
    p_lbl3.font.bold = True
    p_lbl3.font.color.rgb = COLOR_ORANGE
    p_lbl3.line_spacing = 1.05
    
    # Dept details
    dept_box = slide.shapes.add_textbox(Inches(8.6), Inches(5.65), Inches(3.6), Inches(0.7))
    tf_dept = dept_box.text_frame
    tf_dept.word_wrap = True
    tf_dept.margin_left = tf_dept.margin_right = tf_dept.margin_top = tf_dept.margin_bottom = 0
    p_dept = tf_dept.paragraphs[0]
    p_dept.text = "Marian Engineering College\nKazhakuttam\nThiruvananthapuram"
    p_dept.font.name = 'Inter'
    p_dept.font.size = Pt(11)
    p_dept.font.color.rgb = COLOR_WHITE
    p_dept.line_spacing = 1.3

    # Save
    output_filename = "presentation_slide1.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    prs.save(output_path)
    print(f"Slide 1 presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_slide_1()
