import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_BG_DARK = RGBColor(16, 18, 20)      # Dark Card Background (#101214)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(120, 120, 125)    # Muted Gray-green (#78787D)
COLOR_BORDER_THIN = RGBColor(220, 222, 218) # Thin card borders (#DCdeda)
COLOR_WHITE = RGBColor(255, 255, 255)

def set_slide_bg(slide):
    # Set fallback solid color first
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT
    
    # Place background image if available
    image_path = r"C:\Users\gokul_oae3l3i\Downloads\background.png"
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    else:
        print(f"[!] Info: Background image not found at '{image_path}', using solid fallback.")

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def create_slide_8():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 1. Background image setup
    set_slide_bg(slide)
    
    # 2. Slide Index "09" (matching the reference image exactly)
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "09"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(22)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Vertical divider on top left
    draw_line(slide, 1.45, 0.42, 1.46, 0.78, COLOR_BORDER_THIN)
    
    # Top left subheader
    sh_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.4), Inches(3.0), Inches(0.5))
    tf_sh = sh_box.text_frame
    tf_sh.word_wrap = True
    p_sh = tf_sh.paragraphs[0]
    p_sh.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_sh.font.name = 'Inter'
    p_sh.font.size = Pt(8.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_MUTED
    p_sh.line_spacing = 1.15
    
    # 3. Right Top Section Header: CONCLUSION
    draw_line(slide, 11.2, 0.42, 11.215, 0.78, COLOR_ORANGE)
    sol_hdr_box = slide.shapes.add_textbox(Inches(11.35), Inches(0.42), Inches(1.5), Inches(0.4))
    p_shdr = sol_hdr_box.text_frame.paragraphs[0]
    p_shdr.text = "CONCLUSION"
    p_shdr.font.name = 'Inter'
    p_shdr.font.size = Pt(10)
    p_shdr.font.bold = True
    p_shdr.font.color.rgb = COLOR_CHARCOAL
    
    # 4. Main Center Title Box: CONCLUSION
    title_box = slide.shapes.add_textbox(Inches(3.0), Inches(1.9), Inches(7.333), Inches(1.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_h = tf_title.paragraphs[0]
    p_h.alignment = PP_ALIGN.CENTER
    
    r_c = p_h.add_run()
    r_c.text = "CONCLU"
    r_c.font.name = 'Sora'
    r_c.font.size = Pt(48)
    r_c.font.bold = True
    r_c.font.color.rgb = COLOR_CHARCOAL
    
    r_s = p_h.add_run()
    r_s.text = "SION"
    r_s.font.name = 'Sora'
    r_s.font.size = Pt(48)
    r_s.font.bold = True
    r_s.font.color.rgb = COLOR_ORANGE
    
    # Horizontal line under title
    draw_line(slide, 6.1, 3.6, 7.2, 3.62, COLOR_ORANGE)
    
    # 5. Centered Subtitle Paragraph
    sub_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.0), Inches(8.333), Inches(1.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Our vision is to build a future where intelligence\nis not isolated within organizations, but connected\nacross them to create a more collaborative,\nsecure, and intelligent world."
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(13.5)
    p_sub.font.color.rgb = COLOR_CHARCOAL
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.line_spacing = 1.35
    
    # 6. Bottom Quote Highlight Box (Black Background, Centered)
    quote_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.8), Inches(5.5), Inches(9.733), Inches(1.5)
    )
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = COLOR_BG_DARK
    quote_card.line.fill.background()
    
    # Quote Text Box inside card
    quote_box = slide.shapes.add_textbox(Inches(2.1), Inches(5.65), Inches(9.133), Inches(0.7))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = 0
    p_q = tf_q.paragraphs[0]
    p_q.alignment = PP_ALIGN.CENTER
    
    # Orange opening quote
    r_qo = p_q.add_run()
    r_qo.text = "“   "
    r_qo.font.name = 'Libre Baskerville'
    r_qo.font.size = Pt(20)
    r_qo.font.bold = True
    r_qo.font.color.rgb = COLOR_ORANGE
    
    # White body quote text
    r_qt = p_q.add_run()
    r_qt.text = "The future of Artificial Intelligence is not defined by the intelligence\nof individual systems, but by their ability to "
    r_qt.font.name = 'Inter'
    r_qt.font.size = Pt(12)
    r_qt.font.color.rgb = COLOR_WHITE
    
    # Orange keyword
    r_qth = p_q.add_run()
    r_qth.text = "collaborate."
    r_qth.font.name = 'Inter'
    r_qth.font.size = Pt(12)
    r_qth.font.bold = True
    r_qth.font.color.rgb = COLOR_ORANGE
    
    # Orange closing quote
    r_qc = p_q.add_run()
    r_qc.text = "   ”"
    r_qc.font.name = 'Libre Baskerville'
    r_qc.font.size = Pt(20)
    r_qc.font.bold = True
    r_qc.font.color.rgb = COLOR_ORANGE
    
    # Horizontal divider line inside quote card
    draw_line(slide, 6.26, 6.42, 7.06, 6.435, COLOR_ORANGE)
    
    # Footer text below divider inside quote card
    footer_box = slide.shapes.add_textbox(Inches(4.0), Inches(6.52), Inches(5.333), Inches(0.35))
    tf_f = footer_box.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
    p_f = tf_f.paragraphs[0]
    p_f.text = "CONNECT. COLLABORATE. CREATE IMPACT."
    p_f.font.name = 'Inter'
    p_f.font.size = Pt(9.5)
    p_f.font.bold = True
    p_f.font.color.rgb = COLOR_ORANGE
    p_f.alignment = PP_ALIGN.CENTER

    # Save
    output_filename = "presentation_slide8_final.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    try:
        prs.save(output_path)
        print(f"Slide 8 presentation successfully generated at: {output_path}")
    except PermissionError:
        fallback_path = os.path.join(output_dir, "presentation_slide8_fixed.pptx")
        prs.save(fallback_path)
        print(f"[!] WARNING: '{output_path}' was locked/open. Saved fallback to: {fallback_path}")

if __name__ == "__main__":
    create_slide_8()
