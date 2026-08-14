import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BG_LIGHT = RGBColor(246, 244, 239)  # Light Warm Sand (#F6F4EF)
COLOR_BG_DARK = RGBColor(10, 11, 12)     # Very Dark Charcoal/Black (#0A0B0C)
COLOR_CARD = RGBColor(22, 24, 26)        # Dark Card Background (#16181A)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_WHITE = RGBColor(245, 245, 247)    # Off-white text (#F5F5F7)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(92, 95, 92)       # Muted Gray-green (#5C5F5C)

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_lock_icon(slide, cx, cy):
    # Shackle (arch)
    shackle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - 0.05), Inches(cy - 0.08), Inches(0.1), Inches(0.1)
    )
    shackle.fill.background()
    shackle.line.color.rgb = COLOR_ORANGE
    shackle.line.width = Pt(1.5)
    
    # Torso/Body
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - 0.07), Inches(cy - 0.02), Inches(0.14), Inches(0.11)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = COLOR_ORANGE
    body.line.fill.background()
    
    # Keyhole dot
    draw_oval(slide, cx, cy + 0.02, 0.02, COLOR_BG_DARK)

def draw_shield_icon(slide, cx, cy):
    # Custom shield outline using freeform path vertices
    w = 0.08
    h = 0.09
    builder = slide.shapes.build_freeform(Inches(cx - w), Inches(cy - h))
    vertices = [
        (Inches(cx + w), Inches(cy - h)),
        (Inches(cx + w), Inches(cy)),
        (Inches(cx), Inches(cy + h)),
        (Inches(cx - w), Inches(cy)),
        (Inches(cx - w), Inches(cy - h))
    ]
    builder.add_line_segments(vertices, close=True)
    shield = builder.convert_to_shape()
    shield.fill.background()
    shield.line.color.rgb = COLOR_ORANGE
    shield.line.width = Pt(1.5)
    
    # Little inner lock
    draw_oval(slide, cx, cy - 0.02, 0.025, COLOR_ORANGE)
    body = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.03), Inches(cy), Inches(0.06), Inches(0.05)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = COLOR_ORANGE
    body.line.fill.background()

def draw_ecosystem_icon(slide, cx, cy):
    # Draw three small circular nodes
    nodes = [(cx, cy - 0.06), (cx - 0.06, cy + 0.04), (cx + 0.06, cy + 0.04)]
    
    # Connect them with lines
    draw_line(slide, cx - 0.06, cy + 0.04, cx, cy - 0.06, COLOR_ORANGE, 1)
    draw_line(slide, cx + 0.06, cy + 0.04, cx, cy - 0.06, COLOR_ORANGE, 1)
    draw_line(slide, cx - 0.06, cy + 0.04, cx + 0.06, cy + 0.04, COLOR_ORANGE, 1)
    
    # Render node ovals
    for nx, ny in nodes:
        draw_oval(slide, nx, ny, 0.035, COLOR_ORANGE)

def create_slide_2():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Set default slide background to Warm Sand
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT
    
    # Create the left diagonal polygon
    builder = slide.shapes.build_freeform(Inches(0), Inches(0))
    vertices = [
        (Inches(9.0), Inches(0)),
        (Inches(7.7), Inches(7.5)),
        (Inches(0), Inches(7.5)),
        (Inches(0), Inches(0))
    ]
    builder.add_line_segments(vertices, close=True)
    polygon = builder.convert_to_shape()
    polygon.fill.solid()
    polygon.fill.fore_color.rgb = COLOR_BG_DARK
    polygon.line.fill.background()
    
    # Slide Number "02"
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    tf_num.word_wrap = True
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    p_num = tf_num.paragraphs[0]
    p_num.text = "02"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(18)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Vertical Line
    draw_line(slide, 0.55, 1.1, 0.565, 2.5, COLOR_ORANGE)
    
    # Abstract Title
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(6.0), Inches(1.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = "ABSTRACT"
    p_title.font.name = 'Sora'
    p_title.font.size = Pt(58)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    
    # Horizontal divider
    draw_line(slide, 1.1, 3.5, 1.5, 3.53, COLOR_ORANGE)
    
    # Abstract Paragraphs
    abs_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.9), Inches(6.1), Inches(3.2))
    tf_abs = abs_box.text_frame
    tf_abs.word_wrap = True
    tf_abs.margin_left = tf_abs.margin_right = tf_abs.margin_top = tf_abs.margin_bottom = 0
    
    p_abs1 = tf_abs.paragraphs[0]
    p_abs1.text = "As Artificial Intelligence systems become more powerful and widespread, they remain largely isolated within organizational boundaries. This isolation limits collaboration, causes repeated development of similar solutions, and prevents systems from sharing capabilities or knowledge."
    p_abs1.font.name = 'Inter'
    p_abs1.font.size = Pt(13.5)
    p_abs1.font.color.rgb = COLOR_WHITE
    p_abs1.line_spacing = 1.35
    p_abs1.space_after = Pt(24)
    
    p_abs2 = tf_abs.add_paragraph()
    p_abs2.text = "This project proposes a distributed intelligence infrastructure that enables independent AI systems to communicate, collaborate, and coordinate securely while preserving each organization’s autonomy and data privacy."
    p_abs2.font.name = 'Inter'
    p_abs2.font.size = Pt(13.5)
    p_abs2.font.color.rgb = COLOR_WHITE
    p_abs2.line_spacing = 1.35
    
    # Decorative corner barcode wedge lines
    for i in range(12):
        lx = 0.08 + i * 0.08
        lh = 1.2 - i * 0.09
        if lh > 0:
            draw_line(slide, lx, 7.5 - lh, lx + 0.005, 7.5, COLOR_ORANGE)
            
    # Right Column: 3 Feature Blocks
    blocks_data = [
        ("SECURE\nCOLLABORATION", "Enables trusted communication between independent AI systems across organizations.", 1.8, 'shield'),
        ("PRIVACY\nPRESERVATION", "Protects data ownership and ensures privacy through decentralized and secure interaction.", 3.6, 'lock'),
        ("CONNECTED\nAI ECOSYSTEM", "Builds a collaborative ecosystem where AI systems can share capabilities, coordinate tasks, and create value together.", 5.4, 'ecosystem')
    ]
    
    for title, desc, top, icon_type in blocks_data:
        # 1. Circle base
        draw_oval(slide, 9.4, top + 0.35, 0.35, COLOR_BG_DARK)
        
        # 2. Render internal icon
        if icon_type == 'shield':
            draw_shield_icon(slide, 9.4, top + 0.35)
        elif icon_type == 'lock':
            draw_lock_icon(slide, 9.4, top + 0.35)
        elif icon_type == 'ecosystem':
            draw_ecosystem_icon(slide, 9.4, top + 0.35)
            
        # 3. Vertical orange separator line
        draw_line(slide, 10.1, top + 0.05, 10.115, top + 0.65, COLOR_ORANGE)
        
        # 4. Text block next to icon
        text_box = slide.shapes.add_textbox(Inches(10.3), Inches(top), Inches(2.6), Inches(1.3))
        tf_blk = text_box.text_frame
        tf_blk.word_wrap = True
        tf_blk.margin_left = tf_blk.margin_right = tf_blk.margin_top = tf_blk.margin_bottom = 0
        
        p_bt = tf_blk.paragraphs[0]
        p_bt.text = title
        p_bt.font.name = 'Inter'
        p_bt.font.size = Pt(11)
        p_bt.font.bold = True
        p_bt.font.color.rgb = COLOR_CHARCOAL
        p_bt.line_spacing = 1.15
        p_bt.space_after = Pt(8)
        
        p_bd = tf_blk.add_paragraph()
        p_bd.text = desc
        p_bd.font.name = 'Inter'
        p_bd.font.size = Pt(9.5)
        p_bd.font.color.rgb = COLOR_MUTED
        p_bd.line_spacing = 1.2

    # Save
    output_filename = "presentation_slide2.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    prs.save(output_path)
    print(f"Slide 2 presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_slide_2()
