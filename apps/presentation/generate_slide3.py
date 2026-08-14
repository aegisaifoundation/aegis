import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_SAND = RGBColor(246, 244, 239)       # Warm Sand (#F6F4EF)
COLOR_FOREST = RGBColor(31, 77, 58)        # Deep Forest Green (#1F4D3A)
COLOR_ORANGE = RGBColor(198, 106, 43)      # Burnt Orange (#C66A2B)
COLOR_CHARCOAL = RGBColor(30, 30, 30)      # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(92, 95, 92)         # Muted Gray-green (#5C5F5C)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER_SAND = RGBColor(230, 228, 222)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_card_tab(slide, x, y, num_str):
    # Skewed Forest Green tab
    builder = slide.shapes.build_freeform(Inches(x), Inches(y))
    vertices = [
        (Inches(x + 0.75), Inches(y)),
        (Inches(x + 0.55), Inches(y + 0.42)),
        (Inches(x), Inches(y + 0.42)),
        (Inches(x), Inches(y))
    ]
    builder.add_line_segments(vertices, close=True)
    tab = builder.convert_to_shape()
    tab.fill.solid()
    tab.fill.fore_color.rgb = COLOR_FOREST
    tab.line.fill.background()
    
    # Text box for number
    num_box = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.04), Inches(0.4), Inches(0.35))
    tf = num_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = num_str
    p.font.name = 'Inter'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

def draw_user_icon(slide, cx, cy):
    draw_oval(slide, cx, cy - 0.05, 0.04, COLOR_FOREST)
    torso = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - 0.08), Inches(cy + 0.02), Inches(0.16), Inches(0.08)
    )
    torso.fill.solid()
    torso.fill.fore_color.rgb = COLOR_FOREST
    torso.line.fill.background()

def draw_two_users_icon(slide, cx, cy):
    # User 1 (Left background)
    draw_oval(slide, cx - 0.04, cy - 0.05, 0.035, COLOR_MUTED)
    torso1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - 0.11), Inches(cy + 0.01), Inches(0.14), Inches(0.07)
    )
    torso1.fill.solid()
    torso1.fill.fore_color.rgb = COLOR_MUTED
    torso1.line.fill.background()
    
    # User 2 (Right foreground)
    draw_oval(slide, cx + 0.04, cy - 0.03, 0.035, COLOR_FOREST)
    torso2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - 0.03), Inches(cy + 0.03), Inches(0.14), Inches(0.07)
    )
    torso2.fill.solid()
    torso2.fill.fore_color.rgb = COLOR_FOREST
    torso2.line.fill.background()

def draw_doc_icon(slide, cx, cy):
    doc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.07), Inches(cy - 0.09), Inches(0.14), Inches(0.18)
    )
    doc.fill.background()
    doc.line.color.rgb = COLOR_FOREST
    doc.line.width = Pt(1.5)
    
    # Inner lines
    draw_line(slide, cx - 0.04, cy - 0.03, cx + 0.04, cy - 0.03, COLOR_FOREST, 0.8)
    draw_line(slide, cx - 0.04, cy + 0.01, cx + 0.04, cy + 0.01, COLOR_FOREST, 0.8)
    draw_line(slide, cx - 0.04, cy + 0.05, cx + 0.01, cy + 0.05, COLOR_FOREST, 0.8)

def draw_lock_icon(slide, cx, cy):
    shackle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - 0.04), Inches(cy - 0.07), Inches(0.08), Inches(0.08)
    )
    shackle.fill.background()
    shackle.line.color.rgb = COLOR_FOREST
    shackle.line.width = Pt(1.5)
    
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(cx - 0.06), Inches(cy - 0.01), Inches(0.12), Inches(0.09)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = COLOR_FOREST
    body.line.fill.background()
    
    draw_oval(slide, cx, cy + 0.02, 0.015, COLOR_SAND)

def create_slide_3():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Left green column band
    draw_line(slide, 0.0, 0.0, 0.35, 7.5, COLOR_FOREST)
    # Bottom orange segment overlay on left edge
    draw_line(slide, 0.0, 5.0, 0.03, 6.2, COLOR_ORANGE)
    
    # 2. Number "03"
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "03"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(18)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Orange line under number
    draw_line(slide, 0.8, 1.1, 1.2, 1.12, COLOR_ORANGE)
    
    # 3. Header title: Problem Statement
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.0), Inches(1.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = "Problem\nStatement"
    p_title.font.name = 'Libre Baskerville'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_FOREST
    p_title.line_spacing = 1.05
    
    # Orange divider under title
    draw_line(slide, 0.8, 3.1, 1.3, 3.13, COLOR_ORANGE)
    
    # 4. Top-Right Introduction Text
    intro_box = slide.shapes.add_textbox(Inches(6.0), Inches(1.6), Inches(6.5), Inches(1.2))
    tf_intro = intro_box.text_frame
    tf_intro.word_wrap = True
    tf_intro.margin_left = tf_intro.margin_right = tf_intro.margin_top = tf_intro.margin_bottom = 0
    p_intro = tf_intro.paragraphs[0]
    p_intro.text = "Today, AI systems are developed and deployed within isolated organizational boundaries. This fragmentation creates significant barriers to collaboration, leads to duplicated efforts, and slows down collective progress."
    p_intro.font.name = 'Inter'
    p_intro.font.size = Pt(13.5)
    p_intro.font.color.rgb = COLOR_CHARCOAL
    p_intro.line_spacing = 1.3
    
    # 5. Header text on top right
    header_box = slide.shapes.add_textbox(Inches(9.0), Inches(0.58), Inches(3.3), Inches(0.5))
    tf_h = header_box.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_h.font.name = 'Inter'
    p_h.font.size = Pt(8.5)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_MUTED
    p_h.alignment = PP_ALIGN.RIGHT
    p_h.line_spacing = 1.15
    
    # Orange vertical line next to top right header
    draw_line(slide, 12.5, 0.6, 12.515, 0.95, COLOR_ORANGE)
    
    # 6. Row of 4 Cards
    cards_info = [
        ("01", "INDEPENDENT\nAI SYSTEMS", "AI systems operate in silos with no communication or interoperability.", 'user'),
        ("02", "LIMITED\nCOLLABORATION", "Lack of a common platform prevents sharing of capabilities and knowledge.", 'two_users'),
        ("03", "REPEATED\nDEVELOPMENT", "Organizations build similar solutions independently, wasting time and resources.", 'doc'),
        ("04", "NO COMMON\nINFRASTRUCTURE", "Absence of a trusted infrastructure hinders secure and scalable collaboration.", 'lock')
    ]
    
    for i, (num, title, desc, icon_type) in enumerate(cards_info):
        left = 0.8 + i * 3.00
        top = 3.4
        width = 2.7
        height = 2.7
        
        # Draw Warm Sand card background
        card_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_SAND
        card_shape.line.color.rgb = COLOR_BORDER_SAND
        card_shape.line.width = Pt(1)
        
        # Skewed tab for card index number
        draw_card_tab(slide, left, top, num)
        
        # Circle icon outline
        draw_oval(slide, left + 1.35, top + 0.8, 0.275, COLOR_SAND, line_color=COLOR_BORDER_SAND, line_width_pt=1)
        
        # Draw internal icon
        if icon_type == 'user':
            draw_user_icon(slide, left + 1.35, top + 0.8)
        elif icon_type == 'two_users':
            draw_two_users_icon(slide, left + 1.35, top + 0.8)
        elif icon_type == 'doc':
            draw_doc_icon(slide, left + 1.35, top + 0.8)
        elif icon_type == 'lock':
            draw_lock_icon(slide, left + 1.35, top + 0.8)
            
        # Title of card
        tc_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.45), Inches(2.5), Inches(0.6))
        tf_tc = tc_box.text_frame
        tf_tc.word_wrap = True
        p_tc = tf_tc.paragraphs[0]
        p_tc.text = title
        p_tc.font.name = 'Inter'
        p_tc.font.size = Pt(10.5)
        p_tc.font.bold = True
        p_tc.font.color.rgb = COLOR_FOREST
        p_tc.alignment = PP_ALIGN.CENTER
        p_tc.line_spacing = 1.1
        
        # Description of card
        dc_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 2.05), Inches(2.4), Inches(0.6))
        tf_dc = dc_box.text_frame
        tf_dc.word_wrap = True
        tf_dc.margin_left = tf_dc.margin_right = tf_dc.margin_top = tf_dc.margin_bottom = 0
        p_dc = tf_dc.paragraphs[0]
        p_dc.text = desc
        p_dc.font.name = 'Inter'
        p_dc.font.size = Pt(9.5)
        p_dc.font.color.rgb = COLOR_CHARCOAL
        p_dc.alignment = PP_ALIGN.CENTER
        p_dc.line_spacing = 1.25

    # 7. Bottom Quote Highlight Box
    quote_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.8)
    )
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = COLOR_FOREST
    quote_card.line.fill.background()
    
    # Left vertical orange indicator line
    draw_line(slide, 1.2, 6.45, 1.22, 6.95, COLOR_ORANGE)
    
    # Quotes & text box
    # We can write: “ The next challenge... ”
    quote_box = slide.shapes.add_textbox(Inches(1.4), Inches(6.42), Inches(10.5), Inches(0.6))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = 0
    p_q = tf_q.paragraphs[0]
    p_q.text = ""
    
    # Orange opening quote
    r_qo = p_q.add_run()
    r_qo.text = "“   "
    r_qo.font.name = 'Libre Baskerville'
    r_qo.font.size = Pt(20)
    r_qo.font.bold = True
    r_qo.font.color.rgb = COLOR_ORANGE
    
    # White body quote text
    r_qt = p_q.add_run()
    r_qt.text = "The next challenge in AI is enabling independent systems to collaborate securely."
    r_qt.font.name = 'Libre Baskerville'
    r_qt.font.size = Pt(13.5)
    r_qt.font.color.rgb = COLOR_BG_LIGHT
    
    # Orange closing quote
    r_qc = p_q.add_run()
    r_qc.text = "   ”"
    r_qc.font.name = 'Libre Baskerville'
    r_qc.font.size = Pt(20)
    r_qc.font.bold = True
    r_qc.font.color.rgb = COLOR_ORANGE

    # Save
    output_filename = "presentation_slide3.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    prs.save(output_path)
    print(f"Slide 3 presentation successfully generated at: {output_path}")

if __name__ == "__main__":
    create_slide_3()
