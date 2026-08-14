import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_BG_DARK = RGBColor(10, 11, 12)     # Very Dark Charcoal/Black (#0A0B0C)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(120, 120, 125)    # Muted Gray-green (#78787D)
COLOR_BORDER_THIN = RGBColor(220, 222, 218) # Thin card borders (#DCdeda)
COLOR_WHITE = RGBColor(255, 255, 255)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_dashed_vertical_line(slide, cx, y1, y2, color, width_pt=1.0):
    # Standard dashed line via rectangles
    h = 0.05
    gap = 0.04
    curr_y = y1
    while curr_y < y2:
        draw_line(slide, cx - 0.006, curr_y, cx + 0.006, curr_y + h, color)
        curr_y += h + gap

def create_slide_6():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Slide Index "07" (matching the reference image exactly)
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "07"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(22)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Vertical divider on top left
    draw_line(slide, 1.45, 0.42, 1.46, 0.78, COLOR_BORDER_THIN)
    
    # Top left subheader
    sh_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.4), Inches(3.0), Inches(0.5))
    tf_sh = sh_box.text_frame
    tf_sh.word_wrap = True
    p_sh = tf_sh.paragraphs[0]
    p_sh.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_sh.font.name = 'Inter'
    p_sh.font.size = Pt(8.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_MUTED
    p_sh.line_spacing = 1.15
    
    # 2. Right Top Section Header: IMPLEMENTATION PLAN
    draw_line(slide, 10.3, 0.42, 10.315, 0.78, COLOR_ORANGE)
    sol_hdr_box = slide.shapes.add_textbox(Inches(10.45), Inches(0.42), Inches(2.3), Inches(0.4))
    p_shdr = sol_hdr_box.text_frame.paragraphs[0]
    p_shdr.text = "IMPLEMENTATION PLAN"
    p_shdr.font.name = 'Inter'
    p_shdr.font.size = Pt(10)
    p_shdr.font.bold = True
    p_shdr.font.color.rgb = COLOR_CHARCOAL
    
    # 3. Main Center Title Box
    title_box = slide.shapes.add_textbox(Inches(3.0), Inches(1.2), Inches(7.333), Inches(1.1))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_h = tf_title.paragraphs[0]
    p_h.alignment = PP_ALIGN.CENTER
    
    r_i = p_h.add_run()
    r_i.text = "IMPLEMENTATION "
    r_i.font.name = 'Sora'
    r_i.font.size = Pt(44)
    r_i.font.bold = True
    r_i.font.color.rgb = COLOR_CHARCOAL
    
    r_p = p_h.add_run()
    r_p.text = "PLAN"
    r_p.font.name = 'Sora'
    r_p.font.size = Pt(44)
    r_p.font.bold = True
    r_p.font.color.rgb = COLOR_ORANGE
    
    # Horizontal line under title
    draw_line(slide, 6.1, 2.2, 7.2, 2.22, COLOR_ORANGE)
    
    # 4. Horizontal Timeline Line
    draw_line(slide, 1.2, 2.95, 12.133, 2.97, COLOR_CHARCOAL)
    
    # Timeline nodes configurations
    nodes_info = [
        # (cx, num_str, header, bullet_list)
        (1.6, "01", "RESEARCH &\nANALYSIS", [
            "Identify AI collaboration\nchallenges",
            "Study existing approaches",
            "Define project objectives\nand requirements"
        ]),
        (4.15, "02", "SYSTEM\nDESIGN", [
            "Design distributed\narchitecture",
            "Define secure\ncommunication",
            "Plan collaboration\nmechanisms"
        ]),
        (6.7, "03", "PROTOTYPE\nDEVELOPMENT", [
            "Develop a proof\nof concept",
            "Connect multiple\nAI nodes",
            "Demonstrate secure\ncollaboration"
        ]),
        (9.25, "04", "TESTING &\nVALIDATION", [
            "Test communication\nreliability",
            "Evaluate system\nperformance",
            "Verify privacy\nand security"
        ]),
        (11.8, "05", "FUTURE\nEXPANSION", [
            "Support more\norganizations",
            "Improve scalability\nand efficiency",
            "Enable broader\nintelligent collaboration"
        ])
    ]
    
    # Draw timeline nodes and their text columns
    for cx, num_str, header, bullets in nodes_info:
        # Numbered circle (black fill, copper border, diameter = 0.55, radius = 0.275)
        draw_oval(slide, cx, 2.96, 0.275, COLOR_BG_DARK, line_color=COLOR_ORANGE, line_width_pt=1.5)
        
        # Number text inside circle
        num_box = slide.shapes.add_textbox(Inches(cx - 0.2), Inches(2.81), Inches(0.4), Inches(0.3))
        tf_num = num_box.text_frame
        tf_num.word_wrap = True
        tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
        p_num = tf_num.paragraphs[0]
        p_num.text = num_str
        p_num.font.name = 'Inter'
        p_num.font.size = Pt(11.5)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_WHITE
        p_num.alignment = PP_ALIGN.CENTER
        
        # Vertical dashed indicator pointing down
        draw_dashed_vertical_line(slide, cx, 3.32, 3.82, COLOR_ORANGE)
        
        # Column Header Label
        lbl_box = slide.shapes.add_textbox(Inches(cx - 1.1), Inches(4.0), Inches(2.2), Inches(0.6))
        tf_lbl = lbl_box.text_frame
        tf_lbl.word_wrap = True
        tf_lbl.margin_left = tf_lbl.margin_right = tf_lbl.margin_top = tf_lbl.margin_bottom = 0
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = header
        p_lbl.font.name = 'Inter'
        p_lbl.font.size = Pt(10.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_CHARCOAL
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.line_spacing = 1.1
        
        # Tiny orange separator line
        draw_line(slide, cx - 0.2, 4.75, cx + 0.2, 4.765, COLOR_ORANGE)
        
        # Bullet list column
        bullet_box = slide.shapes.add_textbox(Inches(cx - 1.15), Inches(5.0), Inches(2.3), Inches(2.0))
        tf_b = bullet_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
        
        for i, b_text in enumerate(bullets):
            p_b = tf_b.add_paragraph() if tf_b.paragraphs[0].text else tf_b.paragraphs[0]
            p_b.line_spacing = Pt(12)
            p_b.space_after = Pt(8)
            
            # Orange dot run
            r_dot = p_b.add_run()
            r_dot.text = "•   "
            r_dot.font.name = 'Inter'
            r_dot.font.size = Pt(10)
            r_dot.font.bold = True
            r_dot.font.color.rgb = COLOR_ORANGE
            
            # Gray description run
            r_txt = p_b.add_run()
            r_txt.text = b_text
            r_txt.font.name = 'Inter'
            r_txt.font.size = Pt(9.5)
            r_txt.font.color.rgb = COLOR_CHARCOAL

    # Save
    output_filename = "presentation_slide6_final.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    try:
        prs.save(output_path)
        print(f"Slide 6 presentation successfully generated at: {output_path}")
    except PermissionError:
        fallback_path = os.path.join(output_dir, "presentation_slide6_fixed.pptx")
        prs.save(fallback_path)
        print(f"[!] WARNING: '{output_path}' was locked/open. Saved fallback to: {fallback_path}")

if __name__ == "__main__":
    create_slide_6()
