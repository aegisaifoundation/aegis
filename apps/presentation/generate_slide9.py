import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_BG_DARK = RGBColor(10, 11, 12)     # Very Dark Charcoal/Black (#0A0B0C)
COLOR_SAND = RGBColor(246, 244, 239)      # Warm Sand (#F6F4EF)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(120, 120, 125)    # Muted Gray-green (#78787D)
COLOR_BORDER_THIN = RGBColor(220, 222, 218) # Thin card borders (#DCdeda)
COLOR_FOREST = RGBColor(31, 77, 58)        # Deep Forest Green (#1F4D3A)
COLOR_WHITE = RGBColor(255, 255, 255)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

# Tiny outline icons for table headers
def draw_doc_icon(slide, cx, cy, color):
    doc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.06), Inches(cy - 0.08), Inches(0.12), Inches(0.16)
    )
    doc.fill.background()
    doc.line.color.rgb = color
    doc.line.width = Pt(1)
    
    # Inner lines
    draw_line(slide, cx - 0.03, cy - 0.02, cx + 0.03, cy - 0.02, color, 0.8)
    draw_line(slide, cx - 0.03, cy + 0.02, cx + 0.03, cy + 0.02, color, 0.8)

def draw_warning_icon(slide, cx, cy):
    tri = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx - 0.07), Inches(cy - 0.06), Inches(0.14), Inches(0.11)
    )
    tri.fill.background()
    tri.line.color.rgb = COLOR_ORANGE
    tri.line.width = Pt(1)
    # Inner exclamation mark
    draw_line(slide, cx - 0.005, cy - 0.02, cx + 0.005, cy + 0.01, COLOR_ORANGE, 1)
    draw_oval(slide, cx, cy + 0.03, 0.008, COLOR_ORANGE)

def draw_puzzle_icon(slide, cx, cy):
    # Base square
    sq = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.05), Inches(cy - 0.05), Inches(0.1), Inches(0.1)
    )
    sq.fill.background()
    sq.line.color.rgb = COLOR_ORANGE
    sq.line.width = Pt(1)
    
    # Outer connector bumps
    draw_oval(slide, cx + 0.05, cy, 0.025, COLOR_ORANGE)
    draw_oval(slide, cx, cy - 0.05, 0.025, COLOR_ORANGE)

def create_slide_9():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Slide Index Subheader (no slide index circle to match the reference exactly)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.28), Inches(7.0), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_h = tf_title.paragraphs[0]
    
    r_l = p_h.add_run()
    r_l.text = "LITERATURE "
    r_l.font.name = 'Sora'
    r_l.font.size = Pt(36)
    r_l.font.bold = True
    r_l.font.color.rgb = COLOR_CHARCOAL
    
    r_s = p_h.add_run()
    r_s.text = "SURVEY"
    r_s.font.name = 'Sora'
    r_s.font.size = Pt(36)
    r_s.font.bold = True
    r_s.font.color.rgb = COLOR_ORANGE
    
    # Underline
    draw_line(slide, 0.8, 0.95, 1.3, 0.97, COLOR_ORANGE)
    
    # Subheader top right
    sh_box = slide.shapes.add_textbox(Inches(9.0), Inches(0.35), Inches(3.3), Inches(0.5))
    tf_sh = sh_box.text_frame
    tf_sh.word_wrap = True
    p_sh = tf_sh.paragraphs[0]
    p_sh.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_sh.font.name = 'Inter'
    p_sh.font.size = Pt(8.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_MUTED
    p_sh.alignment = PP_ALIGN.RIGHT
    p_sh.line_spacing = 1.15
    
    # Vertical line next to subheader
    draw_line(slide, 12.5, 0.37, 12.515, 0.72, COLOR_ORANGE)
    
    # 2. Custom Table Geometry
    t_top = 1.1
    t_height = 4.95
    row_h = 0.75
    
    col_widths = [3.3, 2.8, 2.8, 2.833]
    col_lefts = [0.8, 4.1, 6.9, 9.7]
    
    # Black Header Background Row
    hdr_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(t_top), Inches(11.733), Inches(0.45)
    )
    hdr_bg.fill.solid()
    hdr_bg.fill.fore_color.rgb = COLOR_BG_DARK
    hdr_bg.line.fill.background()
    
    # Header Columns (Text and Icons)
    headers = [
        ("EXISTING RESEARCH PAPER", 'doc'),
        ("KEY CONTRIBUTION", 'target'),
        ("LIMITATION", 'warning'),
        ("RESEARCH GAP ADDRESSED", 'puzzle')
    ]
    
    for i, (text, icon) in enumerate(headers):
        left = col_lefts[i]
        
        # Icon placement
        icx = left + 0.15
        icy = t_top + 0.225
        
        if icon == 'doc':
            draw_doc_icon(slide, icx, icy, COLOR_ORANGE)
        elif icon == 'target':
            draw_oval(slide, icx, icy, 0.05, COLOR_ORANGE)
            draw_oval(slide, icx, icy, 0.025, COLOR_BG_DARK)
        elif icon == 'warning':
            draw_warning_icon(slide, icx, icy)
        elif icon == 'puzzle':
            draw_puzzle_icon(slide, icx, icy)
            
        # Header Label text
        h_lbl = slide.shapes.add_textbox(Inches(left + 0.35), Inches(t_top + 0.1), Inches(col_widths[i] - 0.4), Inches(0.3))
        p_hl = h_lbl.text_frame.paragraphs[0]
        p_hl.text = text
        p_hl.font.name = 'Inter'
        p_hl.font.size = Pt(8.5)
        p_hl.font.bold = True
        p_hl.font.color.rgb = COLOR_ORANGE
        
    # Table rows content
    rows_data = [
        # Col 1: Title / Author / Venue | Col 2: Key Contribution | Col 3: Limitation | Col 4: Research Gap
        (
            ("Communication-Efficient Learning\nof Deep Networks from\nDecentralized Data", "McMahan et al.", "AISTATS, 2017"),
            "Introduced Federated Learning, enabling collaborative model training without sharing raw data.",
            "Limited to model parameter aggregation; does not support AI-to-AI collaboration or capability sharing.",
            "Need an infrastructure enabling secure collaboration beyond distributed training."
        ),
        (
            ("A Survey on Large\nLanguage Models", "Zhao et al.", "arXiv, 2023"),
            "Comprehensive overview of LLM architectures, capabilities, and applications.",
            "Focuses on individual intelligent models rather than collaboration between independent AI systems.",
            "Enable multiple AI systems to work together across organizational boundaries."
        ),
        (
            ("AutoGen: Enabling Next-Gen\nLLM Applications via\nMulti-Agent Conversation", "Wu et al. (Microsoft)", "arXiv, 2023"),
            "Demonstrated collaborative problem-solving using multiple AI agents through conversational interaction.",
            "Collaboration is confined to a single application and lacks decentralized organizational support.",
            "Support collaboration between autonomous AI systems owned by different organizations."
        ),
        (
            ("Model Context Protocol", "Anthropic", "2024"),
            "Introduced a standardized protocol for connecting AI models with external tools and resources.",
            "Designed for model-tool interaction rather than communication among independent AI infrastructures.",
            "Extend standardized communication to AI-to-AI collaboration across organizations."
        ),
        (
            ("The Rise and Potential of Large\nLanguage Model Based Agents:\nA Survey", "Xi et al.", "arXiv, 2023"),
            "Reviews autonomous AI agents, memory, planning, reasoning, and tool usage.",
            "Primarily focuses on standalone intelligent agents without large-scale inter-organizational collaboration.",
            "Build an infrastructure where autonomous agents can securely discover and collaborate with each other."
        ),
        (
            ("Multi-Agent Reinforcement Learning:\nA Selective Overview of Theories\nand Algorithms", "Zhang, Yang & Başar", "Handbook of Reinforcement Learning and Control, 2021"),
            "Presents theories for coordination, cooperation, and communication among multiple agents.",
            "Mainly theoretical and simulation-oriented, with limited support for enterprise-scale distributed AI collaboration.",
            "Apply multi-agent coordination principles to real-world organizational AI ecosystems."
        )
    ]
    
    for row_idx, data in enumerate(rows_data):
        curr_top = t_top + 0.45 + row_idx * row_h
        
        # Draw horizontal row line at the bottom
        draw_line(slide, 0.8, curr_top + row_h, 12.533, curr_top + row_h, COLOR_BORDER_THIN)
        
        # Col 1: Paper Title & Authors & Icon
        col1_data = data[0]
        # Tiny orange document icon on the left
        draw_doc_icon(slide, col_lefts[0] + 0.15, curr_top + 0.18, COLOR_ORANGE)
        
        p1_box = slide.shapes.add_textbox(Inches(col_lefts[0] + 0.35), Inches(curr_top + 0.04), Inches(col_widths[0] - 0.4), Inches(row_h - 0.08))
        tf1 = p1_box.text_frame
        tf1.word_wrap = True
        tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
        
        p1_t = tf1.paragraphs[0]
        p1_t.text = col1_data[0]
        p1_t.font.name = 'Inter'
        p1_t.font.size = Pt(9.5)
        p1_t.font.bold = True
        p1_t.font.color.rgb = COLOR_CHARCOAL
        p1_t.line_spacing = Pt(11)
        
        p1_a = tf1.add_paragraph()
        p1_a.text = col1_data[1]
        p1_a.font.name = 'Inter'
        p1_a.font.size = Pt(8.5)
        p1_a.font.color.rgb = COLOR_CHARCOAL
        
        p1_v = tf1.add_paragraph()
        p1_v.text = col1_data[2]
        p1_v.font.name = 'Inter'
        p1_v.font.size = Pt(7.5)
        p1_v.font.color.rgb = COLOR_MUTED
        p1_v.font.italic = True
        
        # Col 2: Contribution
        p2_box = slide.shapes.add_textbox(Inches(col_lefts[1] + 0.1), Inches(curr_top + 0.04), Inches(col_widths[1] - 0.2), Inches(row_h - 0.08))
        tf2 = p2_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = data[1]
        p2.font.name = 'Inter'
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = COLOR_CHARCOAL
        p2.line_spacing = Pt(11.5)
        
        # Col 3: Limitation
        p3_box = slide.shapes.add_textbox(Inches(col_lefts[2] + 0.1), Inches(curr_top + 0.04), Inches(col_widths[2] - 0.2), Inches(row_h - 0.08))
        tf3 = p3_box.text_frame
        tf3.word_wrap = True
        tf3.margin_left = tf3.margin_right = tf3.margin_top = tf3.margin_bottom = 0
        p3 = tf3.paragraphs[0]
        p3.text = data[2]
        p3.font.name = 'Inter'
        p3.font.size = Pt(8.5)
        p3.font.color.rgb = COLOR_CHARCOAL
        p3.line_spacing = Pt(11.5)
        
        # Col 4: Research Gap
        p4_box = slide.shapes.add_textbox(Inches(col_lefts[3] + 0.1), Inches(curr_top + 0.04), Inches(col_widths[3] - 0.2), Inches(row_h - 0.08))
        tf4 = p4_box.text_frame
        tf4.word_wrap = True
        tf4.margin_left = tf4.margin_right = tf4.margin_top = tf4.margin_bottom = 0
        p4 = tf4.paragraphs[0]
        p4.text = data[3]
        p4.font.name = 'Inter'
        p4.font.size = Pt(8.5)
        p4.font.color.rgb = COLOR_CHARCOAL
        p4.line_spacing = Pt(11.5)
        
    # Draw vertical column lines (3 lines)
    for lx in [4.1, 6.9, 9.7]:
        draw_line(slide, lx, t_top, lx + 0.005, t_top + t_height, COLOR_BORDER_THIN)
        
    # 3. Bottom Research Gap Highlight Box
    quote_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(6.22), Inches(11.733), Inches(0.95)
    )
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = COLOR_SAND
    quote_card.line.color.rgb = COLOR_BORDER_THIN
    quote_card.line.width = Pt(1)
    
    # Forest Green circle target icon (Forest Green with Orange target inside)
    draw_oval(slide, 1.3, 6.70, 0.25, COLOR_FOREST)
    draw_oval(slide, 1.3, 6.70, 0.15, COLOR_FOREST, line_color=COLOR_ORANGE, line_width_pt=1.2)
    draw_oval(slide, 1.3, 6.70, 0.05, COLOR_ORANGE)
    
    # "RESEARCH GAP" title text block next to it
    rg_box = slide.shapes.add_textbox(Inches(1.68), Inches(6.5), Inches(1.0), Inches(0.4))
    tf_rg = rg_box.text_frame
    tf_rg.word_wrap = True
    tf_rg.margin_left = tf_rg.margin_right = tf_rg.margin_top = tf_rg.margin_bottom = 0
    p_rg = tf_rg.paragraphs[0]
    p_rg.text = "RESEARCH\nGAP"
    p_rg.font.name = 'Inter'
    p_rg.font.size = Pt(10)
    p_rg.font.bold = True
    p_rg.font.color.rgb = COLOR_FOREST
    p_rg.line_spacing = 1.05
    
    # Left vertical orange indicator line
    draw_line(slide, 2.7, 6.35, 2.715, 7.05, COLOR_ORANGE)
    
    # Body description quote text box
    desc_box = slide.shapes.add_textbox(Inches(2.88), Inches(6.28), Inches(9.5), Inches(0.85))
    tf_d = desc_box.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_right = tf_d.margin_top = tf_d.margin_bottom = 0
    p_d = tf_d.paragraphs[0]
    p_d.line_spacing = Pt(11.2)
    p_d.text = ""
    
    r1 = p_d.add_run()
    r1.text = "Although existing research has significantly advanced Federated Learning, Large Language Models, Multi-Agent Systems, AI Agents, and AI communication protocols, these studies solve only specific aspects of distributed intelligence. Current approaches focus on collaborative model training, agent coordination within controlled environments, or interaction between AI models and external tools. "
    r1.font.name = 'Inter'
    r1.font.size = Pt(8.5)
    r1.font.color.rgb = COLOR_CHARCOAL
    
    r2 = p_d.add_run()
    r2.text = "No existing work provides a unified infrastructure that enables independent AI systems across different organizations to securely discover each other, communicate, share capabilities, coordinate tasks, and collaboratively solve problems while preserving privacy, autonomy, and organizational ownership."
    r2.font.name = 'Inter'
    r2.font.size = Pt(8.5)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_ORANGE
    
    r3 = p_d.add_run()
    r3.text = " This gap motivates the development of a distributed intelligence infrastructure for connecting AI ecosystems."
    r3.font.name = 'Inter'
    r3.font.size = Pt(8.5)
    r3.font.color.rgb = COLOR_CHARCOAL

    # Save
    output_filename = "presentation_slide9_final.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    try:
        prs.save(output_path)
        print(f"Slide 9 presentation successfully generated at: {output_path}")
    except PermissionError:
        fallback_path = os.path.join(output_dir, "presentation_slide9_fixed.pptx")
        prs.save(fallback_path)
        print(f"[!] WARNING: '{output_path}' was locked/open. Saved fallback to: {fallback_path}")

if __name__ == "__main__":
    create_slide_9()
