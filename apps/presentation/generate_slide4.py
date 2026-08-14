import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_BG_DARK = RGBColor(10, 11, 12)     # Very Dark Charcoal/Black (#0A0B0C)
COLOR_CARD = RGBColor(22, 24, 26)        # Dark Card Background (#16181A)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_WHITE = RGBColor(245, 245, 247)    # Off-white text (#F5F5F7)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(140, 140, 145)    # Muted Gray-green (#8C8C91)
COLOR_BORDER_THIN = RGBColor(220, 222, 218) # Thin card borders (#DCdeda)
COLOR_FOREST = RGBColor(31, 77, 58)        # Deep Forest Green (#1F4D3A)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_connector(slide, x1, y1, x2, y2, color, width_pt=1.0):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn

# Helper vector shapes for icons
def draw_hospital_icon(slide, cx, cy, color):
    # Main building
    b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(cx - 0.09), Inches(cy - 0.04), Inches(0.18), Inches(0.12)
    )
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    
    # Roof cross
    draw_line(slide, cx - 0.015, cy - 0.09, cx + 0.015, cy - 0.04, color, 1.5)
    draw_line(slide, cx - 0.04, cy - 0.065, cx + 0.04, cy - 0.065, color, 1.5)

def draw_school_icon(slide, cx, cy, color):
    # Triangle roof
    r = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx - 0.09), Inches(cy - 0.08), Inches(0.18), Inches(0.06)
    )
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()
    
    # Base columns block
    c = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.08), Inches(cy - 0.02), Inches(0.16), Inches(0.08)
    )
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    
    # Tiny door cut
    draw_oval(slide, cx, cy + 0.02, 0.018, COLOR_BG_LIGHT)

def draw_gov_icon(slide, cx, cy, color):
    # Base pillars
    b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.09), Inches(cy - 0.04), Inches(0.18), Inches(0.10)
    )
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()
    
    # Pediment triangle roof
    r = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx - 0.1), Inches(cy - 0.09), Inches(0.20), Inches(0.05)
    )
    r.fill.solid()
    r.fill.fore_color.rgb = color
    r.line.fill.background()

def draw_industry_icon(slide, cx, cy, color):
    # Factory saw-tooth shape
    builder = slide.shapes.build_freeform(Inches(cx - 0.09), Inches(cy + 0.06))
    vertices = [
        (Inches(cx - 0.09), Inches(cy - 0.03)),
        (Inches(cx - 0.04), Inches(cy + 0.01)),
        (Inches(cx - 0.04), Inches(cy - 0.03)),
        (Inches(cx + 0.01), Inches(cy + 0.01)),
        (Inches(cx + 0.01), Inches(cy - 0.03)),
        (Inches(cx + 0.07), Inches(cy + 0.01)),
        (Inches(cx + 0.07), Inches(cy + 0.06)),
        (Inches(cx - 0.09), Inches(cy + 0.06))
    ]
    builder.add_line_segments(vertices, close=True)
    f = builder.convert_to_shape()
    f.fill.solid()
    f.fill.fore_color.rgb = color
    f.line.fill.background()

def draw_badge_shield(slide, cx, cy):
    # Custom shield outline
    w = 0.04
    h = 0.04
    builder = slide.shapes.build_freeform(Inches(cx - w), Inches(cy - h))
    vertices = [
        (Inches(cx + w), Inches(cy - h)),
        (Inches(cx + w), Inches(cy)),
        (Inches(cx), Inches(cy + h)),
        (Inches(cx - w), Inches(cy)),
        (Inches(cx - w), Inches(cy - h))
    ]
    builder.add_line_segments(vertices, close=True)
    shield = builder.convert_to_shape()
    shield.fill.background()
    shield.line.color.rgb = COLOR_ORANGE
    shield.line.width = Pt(1)

def draw_badge_cap(slide, cx, cy):
    # Diamond graduation cap
    cap = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(cx - 0.05), Inches(cy - 0.04), Inches(0.1), Inches(0.06)
    )
    cap.fill.solid()
    cap.fill.fore_color.rgb = COLOR_ORANGE
    cap.line.fill.background()

def draw_badge_building(slide, cx, cy):
    # Tiny pillars
    pillars = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.04), Inches(cy - 0.03), Inches(0.08), Inches(0.06)
    )
    pillars.fill.solid()
    pillars.fill.fore_color.rgb = COLOR_ORANGE
    pillars.line.fill.background()

def draw_badge_factory(slide, cx, cy):
    # Tiny factory box
    draw_oval(slide, cx, cy, 0.04, COLOR_CARD, line_color=COLOR_ORANGE, line_width_pt=1)

def draw_detailed_hub_network(slide, cx, cy):
    # Center node
    draw_oval(slide, cx, cy, 0.025, COLOR_ORANGE)
    # 6 surrounding nodes at radius 0.095 inches
    r_node = 0.095
    nodes = []
    for i in range(6):
        angle = math.radians(i * 60)
        nx = cx + r_node * math.cos(angle)
        ny = cy + r_node * math.sin(angle)
        nodes.append((nx, ny))
        # Draw node circle
        draw_oval(slide, nx, ny, 0.02, COLOR_ORANGE)
        # Connect to center
        draw_connector(slide, cx, cy, nx, ny, COLOR_ORANGE, 0.8)
    # Connect surrounding nodes in a hexagon ring
    for i in range(6):
        n1 = nodes[i]
        n2 = nodes[(i + 1) % 6]
        draw_connector(slide, n1[0], n1[1], n2[0], n2[1], COLOR_ORANGE, 0.8)

# Outlines for Cardinal Text Blocks
def draw_lock_outline(slide, cx, cy):
    shackle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - 0.035), Inches(cy - 0.07), Inches(0.07), Inches(0.07)
    )
    shackle.fill.background()
    shackle.line.color.rgb = COLOR_ORANGE
    shackle.line.width = Pt(1)
    
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - 0.05), Inches(cy - 0.02), Inches(0.1), Inches(0.08)
    )
    body.fill.background()
    body.line.color.rgb = COLOR_ORANGE
    body.line.width = Pt(1)

def draw_db_stack(slide, cx, cy):
    for dy in [-0.05, 0, 0.05]:
        c = draw_oval(slide, cx, cy + dy, 0.05, COLOR_BG_LIGHT, line_color=COLOR_ORANGE, line_width_pt=1)
        c.fill.background()

def draw_brain_outline(slide, cx, cy):
    nodes = [(cx, cy - 0.04), (cx - 0.04, cy + 0.03), (cx + 0.04, cy + 0.03)]
    draw_connector(slide, cx - 0.04, cy + 0.03, cx, cy - 0.04, COLOR_ORANGE, 0.8)
    draw_connector(slide, cx + 0.04, cy + 0.03, cx, cy - 0.04, COLOR_ORANGE, 0.8)
    draw_connector(slide, cx - 0.04, cy + 0.03, cx + 0.04, cy + 0.03, COLOR_ORANGE, 0.8)
    for nx, ny in nodes:
        draw_oval(slide, nx, ny, 0.02, COLOR_BG_LIGHT, line_color=COLOR_ORANGE, line_width_pt=1)

def draw_check_outline(slide, cx, cy):
    draw_oval(slide, cx, cy, 0.045, COLOR_BG_LIGHT, line_color=COLOR_ORANGE, line_width_pt=1)
    draw_line(slide, cx - 0.018, cy, cx - 0.005, cy + 0.015, COLOR_ORANGE, 1)
    draw_line(slide, cx - 0.005, cy + 0.015, cx + 0.018, cy - 0.015, COLOR_ORANGE, 1)

def create_slide_4():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Slide Index "04"
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "04"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(22)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Vertical divider on top left
    draw_line(slide, 1.45, 0.42, 1.46, 0.78, COLOR_BORDER_THIN)
    
    # Top left subheader
    sh_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.4), Inches(3.0), Inches(0.5))
    tf_sh = sh_box.text_frame
    tf_sh.word_wrap = True
    p_sh = tf_sh.paragraphs[0]
    p_sh.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_sh.font.name = 'Inter'
    p_sh.font.size = Pt(8.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_MUTED
    p_sh.line_spacing = 1.15
    
    # 2. Right Top Section Header: OUR SOLUTION
    draw_line(slide, 10.3, 0.42, 10.315, 0.78, COLOR_ORANGE)
    sol_hdr_box = slide.shapes.add_textbox(Inches(10.45), Inches(0.42), Inches(2.0), Inches(0.4))
    p_shdr = sol_hdr_box.text_frame.paragraphs[0]
    p_shdr.text = "OUR SOLUTION"
    p_shdr.font.name = 'Inter'
    p_shdr.font.size = Pt(10)
    p_shdr.font.bold = True
    p_shdr.font.color.rgb = COLOR_CHARCOAL
    
    # 3. Left Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(4.5), Inches(2.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    
    p_lbl = tf_title.paragraphs[0]
    p_lbl.text = "THE SOLUTION"
    p_lbl.font.name = 'Inter'
    p_lbl.font.size = Pt(10)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = COLOR_ORANGE
    p_lbl.space_after = Pt(8)
    
    # Sora heading runs
    p_h = tf_title.add_paragraph()
    p_h.line_spacing = 0.95
    p_h.space_after = Pt(8)
    
    r_d = p_h.add_run()
    r_d.text = "DISTRIBUTED\n"
    r_d.font.name = 'Sora'
    r_d.font.size = Pt(36)
    r_d.font.bold = True
    r_d.font.color.rgb = COLOR_CHARCOAL
    
    r_i = p_h.add_run()
    r_i.text = "INTELLIGENCE\n"
    r_i.font.name = 'Sora'
    r_i.font.size = Pt(36)
    r_i.font.bold = True
    r_i.font.color.rgb = COLOR_ORANGE
    
    r_inf = p_h.add_run()
    r_inf.text = "INFRASTRUCTURE"
    r_inf.font.name = 'Sora'
    r_inf.font.size = Pt(36)
    r_inf.font.bold = True
    r_inf.font.color.rgb = COLOR_CHARCOAL
    
    # Horizontal line under title
    draw_line(slide, 0.8, 3.5, 1.3, 3.52, COLOR_ORANGE)
    
    # Description paragraph
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.85), Inches(3.6), Inches(1.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = 0
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A secure, decentralized foundation that enables independent AI systems to connect, collaborate, and create collective impact."
    p_desc.font.name = 'Inter'
    p_desc.font.size = Pt(12.5)
    p_desc.font.color.rgb = COLOR_CHARCOAL
    p_desc.line_spacing = 1.35
    
    # Target Box (Bottom Left)
    tbox_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(5.15), Inches(4.0), Inches(0.9)
    )
    tbox_bg.fill.solid()
    tbox_bg.fill.fore_color.rgb = COLOR_BG_DARK
    tbox_bg.line.fill.background()
    
    # Target Icon Outline Circle
    draw_oval(slide, 1.3, 5.6, 0.22, COLOR_BG_DARK, line_color=COLOR_ORANGE, line_width_pt=1.2)
    # Simple target symbol
    draw_oval(slide, 1.3, 5.6, 0.12, COLOR_ORANGE)
    draw_oval(slide, 1.3, 5.6, 0.05, COLOR_BG_DARK)
    
    # Target box text
    tbox_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.27), Inches(2.9), Inches(0.75))
    tf_tb = tbox_text.text_frame
    tf_tb.word_wrap = True
    p_tb1 = tf_tb.paragraphs[0]
    p_tb1.text = "Connecting independent intelligence."
    p_tb1.font.name = 'Inter'
    p_tb1.font.size = Pt(10.5)
    p_tb1.font.color.rgb = COLOR_WHITE
    p_tb1.space_after = Pt(2)
    
    p_tb2 = tf_tb.add_paragraph()
    p_tb2.text = "Empowering collective progress."
    p_tb2.font.name = 'Inter'
    p_tb2.font.size = Pt(10.5)
    p_tb2.font.bold = True
    p_tb2.font.color.rgb = COLOR_ORANGE
    
    # 4. Right Side Diagram (Central Hub & Diagonal Outer Circles)
    hcx, hcy = 9.0, 3.7
    
    # Outer Dotted Circle Orbit (radius = 1.3 inches)
    draw_oval(slide, hcx, hcy, 1.3, COLOR_BG_LIGHT, line_color=COLOR_BORDER_THIN, line_width_pt=0.8)
    
    # Diagonal outer circles coordinates (radius = 2.2 inches, R_cos = 1.55 inches)
    nodes_info = [
        # (cx, cy, label, icon_type, badge_type)
        (7.45, 2.15, "HEALTHCARE\nORGANIZATION", 'hospital', 'shield'),
        (10.55, 2.15, "EDUCATIONAL\nINSTITUTION", 'school', 'cap'),
        (10.55, 5.25, "GOVERNMENT\nAGENCY", 'gov', 'building'),
        (7.45, 5.25, "INDUSTRY\nENTERPRISE", 'industry', 'factory')
    ]
    
    # Draw connectors to hub centers
    for cx, cy, label, icon_type, badge_type in nodes_info:
        draw_connector(slide, hcx, hcy, cx, cy, COLOR_BORDER_THIN, 1.0)
        
    # Draw the 4 diagonal orange beads ON the dotted orbit (1.3 inches radius)
    # Angle 45, 135, 225, 315
    for dx_sign, dy_sign in [(-1, -1), (1, -1), (1, 1), (-1, 1)]:
        bcx = hcx + dx_sign * 1.3 * 0.707
        bcy = hcy + dy_sign * 1.3 * 0.707
        draw_oval(slide, bcx, bcy, 0.045, COLOR_ORANGE)
        
    # Draw the 4 cardinal orange beads ON the dotted orbit (1.3 inches radius)
    # Angle 0, 90, 180, 270
    draw_oval(slide, hcx, hcy - 1.3, 0.045, COLOR_ORANGE)  # Top (90°)
    draw_oval(slide, hcx + 1.3, hcy, 0.045, COLOR_ORANGE)  # Right (0°)
    draw_oval(slide, hcx, hcy + 1.3, 0.045, COLOR_ORANGE)  # Bottom (270°)
    draw_oval(slide, hcx - 1.3, hcy, 0.045, COLOR_ORANGE)  # Left (180°)
    
    # Draw central hub circle
    draw_oval(slide, hcx, hcy, 0.95, COLOR_BG_DARK, line_color=COLOR_ORANGE, line_width_pt=1.5)
    
    # Central Hub Text & Hexagon Network Mesh Icon
    draw_detailed_hub_network(slide, hcx, hcy - 0.20)
    hub_box = slide.shapes.add_textbox(Inches(hcx - 0.8), Inches(hcy + 0.08), Inches(1.6), Inches(0.6))
    tf_hb = hub_box.text_frame
    tf_hb.word_wrap = True
    tf_hb.margin_left = tf_hb.margin_right = tf_hb.margin_top = tf_hb.margin_bottom = 0
    p_hb = tf_hb.paragraphs[0]
    p_hb.text = "DISTRIBUTED\nINTELLIGENCE\nINFRASTRUCTURE"
    p_hb.font.name = 'Inter'
    p_hb.font.size = Pt(8.5)
    p_hb.font.bold = True
    p_hb.font.color.rgb = COLOR_WHITE
    p_hb.alignment = PP_ALIGN.CENTER
    p_hb.line_spacing = 1.05
    
    # Render outer circles
    for cx, cy, label, icon_type, badge_type in nodes_info:
        # White base circle (diameter = 1.2 inches, radius = 0.6)
        draw_oval(slide, cx, cy, 0.6, COLOR_BG_LIGHT, line_color=COLOR_BORDER_THIN, line_width_pt=1)
        
        # Center icon (drawn in Forest Green, placed in upper half of outer circle)
        if icon_type == 'hospital':
            draw_hospital_icon(slide, cx, cy - 0.12, COLOR_FOREST)
            bcx, bcy = cx - 0.50, cy - 0.30
        elif icon_type == 'school':
            draw_school_icon(slide, cx, cy - 0.12, COLOR_FOREST)
            bcx, bcy = cx + 0.50, cy - 0.30
        elif icon_type == 'gov':
            draw_gov_icon(slide, cx, cy - 0.12, COLOR_FOREST)
            bcx, bcy = cx + 0.50, cy + 0.30
        elif icon_type == 'industry':
            draw_industry_icon(slide, cx, cy - 0.12, COLOR_FOREST)
            bcx, bcy = cx - 0.50, cy + 0.30
            
        # Draw black badge circle (diameter = 0.3, radius = 0.15)
        draw_oval(slide, bcx, bcy, 0.15, COLOR_BG_DARK)
        
        # Render tiny orange badge icon
        if badge_type == 'shield':
            draw_badge_shield(slide, bcx, bcy)
        elif badge_type == 'cap':
            draw_badge_cap(slide, bcx, bcy)
        elif badge_type == 'building':
            draw_badge_building(slide, bcx, bcy)
        elif badge_type == 'factory':
            draw_badge_factory(slide, bcx, bcy)
            
        # Small Label placed inside bottom half of circle
        lbl_box = slide.shapes.add_textbox(Inches(cx - 0.55), Inches(cy + 0.15), Inches(1.1), Inches(0.35))
        tf_lbl = lbl_box.text_frame
        tf_lbl.word_wrap = True
        tf_lbl.margin_left = tf_lbl.margin_right = tf_lbl.margin_top = tf_lbl.margin_bottom = 0
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.name = 'Inter'
        p_lbl.font.size = Pt(7.5)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = COLOR_FOREST
        p_lbl.alignment = PP_ALIGN.CENTER
        p_lbl.line_spacing = 1.05

    # Cardinal Text Blocks & Outlines
    # 1. Top cardinal section (above the top dot)
    draw_lock_outline(slide, 9.0, 1.45)
    top_lbl = slide.shapes.add_textbox(Inches(7.5), Inches(0.50), Inches(3.0), Inches(0.65))
    tf_top = top_lbl.text_frame
    tf_top.word_wrap = True
    p_top = tf_top.paragraphs[0]
    p_top.alignment = PP_ALIGN.CENTER
    r = p_top.add_run()
    r.text = "SECURE BY DESIGN\n"
    r.font.name = 'Inter'
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = COLOR_CHARCOAL
    r2 = p_top.add_run()
    r2.text = "Privacy-preserving collaboration\nwith zero-trust architecture."
    r2.font.name = 'Inter'
    r2.font.size = Pt(8)
    r2.font.color.rgb = COLOR_MUTED
    p_top.line_spacing = 1.15

    # 2. Bottom cardinal section (below the bottom dot)
    draw_check_outline(slide, 9.0, 5.75)
    bot_lbl = slide.shapes.add_textbox(Inches(7.5), Inches(6.0), Inches(3.0), Inches(0.65))
    tf_bot = bot_lbl.text_frame
    tf_bot.word_wrap = True
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.CENTER
    r = p_bot.add_run()
    r.text = "DECENTRALIZED & SOVEREIGN\n"
    r.font.name = 'Inter'
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = COLOR_CHARCOAL
    r2 = p_bot.add_run()
    r2.text = "Each organization retains control\nover its data and models."
    r2.font.name = 'Inter'
    r2.font.size = Pt(8)
    r2.font.color.rgb = COLOR_MUTED
    p_bot.line_spacing = 1.15

    # 3. Left cardinal section (left of the left dot)
    draw_db_stack(slide, 6.75, 3.7)
    left_lbl = slide.shapes.add_textbox(Inches(4.1), Inches(3.3), Inches(2.4), Inches(0.8))
    tf_left = left_lbl.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0
    p_left = tf_left.paragraphs[0]
    p_left.alignment = PP_ALIGN.RIGHT
    r = p_left.add_run()
    r.text = "SHARED CAPABILITIES\n"
    r.font.name = 'Inter'
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = COLOR_CHARCOAL
    r2 = p_left.add_run()
    r2.text = "Discover and access\ncapabilities across\norganizations."
    r2.font.name = 'Inter'
    r2.font.size = Pt(8)
    r2.font.color.rgb = COLOR_MUTED
    p_left.line_spacing = 1.15

    # 4. Right cardinal section (right of the right dot)
    draw_brain_outline(slide, 11.25, 3.7)
    right_lbl = slide.shapes.add_textbox(Inches(11.55), Inches(3.3), Inches(2.4), Inches(0.8))
    tf_right = right_lbl.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0
    p_right = tf_right.paragraphs[0]
    p_right.alignment = PP_ALIGN.LEFT
    r = p_right.add_run()
    r.text = "COLLECTIVE INTELLIGENCE\n"
    r.font.name = 'Inter'
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = COLOR_CHARCOAL
    r2 = p_right.add_run()
    r2.text = "Combine strengths,\nshare insights, and\nsolve bigger challenges."
    r2.font.name = 'Inter'
    r2.font.size = Pt(8)
    r2.font.color.rgb = COLOR_MUTED
    p_right.line_spacing = 1.15

    # 5. Bottom Horizontal Quote Band
    quote_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0), Inches(7.0), Inches(13.333), Inches(0.5)
    )
    quote_band.fill.solid()
    quote_band.fill.fore_color.rgb = COLOR_BG_DARK
    quote_band.line.fill.background()
    
    quote_box = slide.shapes.add_textbox(Inches(1.0), Inches(7.05), Inches(11.333), Inches(0.4))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = 0
    p_q = tf_q.paragraphs[0]
    p_q.alignment = PP_ALIGN.CENTER
    
    r_qo = p_q.add_run()
    r_qo.text = "“   "
    r_qo.font.name = 'Libre Baskerville'
    r_qo.font.size = Pt(16)
    r_qo.font.bold = True
    r_qo.font.color.rgb = COLOR_ORANGE
    
    r_qt = p_q.add_run()
    r_qt.text = "Turning isolated AI systems into a network of "
    r_qt.font.name = 'Libre Baskerville'
    r_qt.font.size = Pt(13)
    r_qt.font.color.rgb = COLOR_WHITE
    
    r_qth = p_q.add_run()
    r_qth.text = "collaborative intelligence"
    r_qth.font.name = 'Libre Baskerville'
    r_qth.font.size = Pt(13)
    r_qth.font.bold = True
    r_qth.font.color.rgb = COLOR_ORANGE
    
    r_qt2 = p_q.add_run()
    r_qt2.text = "."
    r_qt2.font.name = 'Libre Baskerville'
    r_qt2.font.size = Pt(13)
    r_qt2.font.color.rgb = COLOR_WHITE
    
    r_qc = p_q.add_run()
    r_qc.text = "   ”"
    r_qc.font.name = 'Libre Baskerville'
    r_qc.font.size = Pt(16)
    r_qc.font.bold = True
    r_qc.font.color.rgb = COLOR_ORANGE

    # Save
    output_filename = "presentation_slide4_final.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    try:
        prs.save(output_path)
        print(f"Slide 4 presentation successfully generated at: {output_path}")
    except PermissionError:
        fallback_path = os.path.join(output_dir, "presentation_slide4_fixed.pptx")
        prs.save(fallback_path)
        print(f"[!] WARNING: '{output_path}' was locked/open. Saved fallback to: {fallback_path}")

if __name__ == "__main__":
    create_slide_4()
