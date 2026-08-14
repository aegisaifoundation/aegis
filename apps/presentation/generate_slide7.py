import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# Colors
COLOR_BG_LIGHT = RGBColor(252, 252, 250)  # Off White (#FCFCFA)
COLOR_BG_DARK = RGBColor(10, 11, 12)     # Very Dark Charcoal/Black (#0A0B0C)
COLOR_SAND = RGBColor(246, 244, 239)      # Warm Sand (#F6F4EF)
COLOR_ORANGE = RGBColor(198, 106, 43)    # Orange/Copper accent (#C66A2B)
COLOR_CHARCOAL = RGBColor(30, 30, 30)    # Charcoal text (#1E1E1E)
COLOR_MUTED = RGBColor(120, 120, 125)    # Muted Gray-green (#78787D)
COLOR_BORDER_THIN = RGBColor(220, 222, 218) # Thin card borders (#DCdeda)
COLOR_FOREST = RGBColor(31, 77, 58)        # Deep Forest Green (#1F4D3A)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(x1), Inches(y1), Inches(x2 - x1 if (x2 - x1) > 0 else 0.015), Inches(y2 - y1 if (y2 - y1) > 0 else 0.015)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def draw_connector(slide, x1, y1, x2, y2, color, width_pt=1.0):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn

# Helper vector shapes for icons
def draw_icon_shield_check(slide, cx, cy):
    # Shield outline
    w = 0.06
    h = 0.07
    builder = slide.shapes.build_freeform(Inches(cx - w), Inches(cy - h))
    vertices = [
        (Inches(cx + w), Inches(cy - h)),
        (Inches(cx + w), Inches(cy)),
        (Inches(cx), Inches(cy + h)),
        (Inches(cx - w), Inches(cy)),
        (Inches(cx - w), Inches(cy - h))
    ]
    builder.add_line_segments(vertices, close=True)
    shield = builder.convert_to_shape()
    shield.fill.background()
    shield.line.color.rgb = COLOR_FOREST
    shield.line.width = Pt(1)
    # Checkmark inside
    draw_line(slide, cx - 0.02, cy, cx - 0.005, cy + 0.015, COLOR_FOREST, 0.8)
    draw_line(slide, cx - 0.005, cy + 0.015, cx + 0.02, cy - 0.015, COLOR_FOREST, 0.8)

def draw_icon_lightbulb(slide, cx, cy):
    # Bulb glass outline
    draw_oval(slide, cx, cy - 0.02, 0.065, COLOR_BG_LIGHT, line_color=COLOR_FOREST, line_width_pt=1).fill.background()
    # Base filament connector
    draw_line(slide, cx - 0.035, cy + 0.05, cx + 0.035, cy + 0.06, COLOR_FOREST)
    draw_line(slide, cx - 0.02, cy + 0.08, cx + 0.02, cy + 0.09, COLOR_FOREST)
    # Light ray lines
    for angle in [0, 45, 90, 135, 180, 225, 270, 315]:
        rad = math.radians(angle)
        x1 = cx + 0.08 * math.cos(rad)
        y1 = cy - 0.02 + 0.08 * math.sin(rad)
        x2 = cx + 0.12 * math.cos(rad)
        y2 = cy - 0.02 + 0.12 * math.sin(rad)
        draw_connector(slide, x1, y1, x2, y2, COLOR_FOREST, 0.8)

def draw_icon_group(slide, cx, cy):
    # Three user outlines
    # User 1 (Center)
    draw_oval(slide, cx, cy - 0.04, 0.03, COLOR_FOREST)
    torso1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - 0.06), Inches(cy + 0.01), Inches(0.12), Inches(0.06)
    )
    torso1.fill.solid()
    torso1.fill.fore_color.rgb = COLOR_FOREST
    torso1.line.fill.background()
    
    # User 2 (Left)
    draw_oval(slide, cx - 0.06, cy + 0.01, 0.025, COLOR_MUTED)
    torso2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - 0.11), Inches(cy + 0.05), Inches(0.10), Inches(0.05)
    )
    torso2.fill.solid()
    torso2.fill.fore_color.rgb = COLOR_MUTED
    torso2.line.fill.background()
    
    # User 3 (Right)
    draw_oval(slide, cx + 0.06, cy + 0.01, 0.025, COLOR_MUTED)
    torso3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx + 0.01), Inches(cy + 0.05), Inches(0.10), Inches(0.05)
    )
    torso3.fill.solid()
    torso3.fill.fore_color.rgb = COLOR_MUTED
    torso3.line.fill.background()

def draw_icon_chart(slide, cx, cy):
    # 3 chart pillars
    draw_line(slide, cx - 0.07, cy + 0.03, cx - 0.04, cy + 0.09, COLOR_FOREST)
    draw_line(slide, cx - 0.01, cy - 0.01, cx + 0.02, cy + 0.09, COLOR_FOREST)
    draw_line(slide, cx + 0.05, cy - 0.05, cx + 0.08, cy + 0.09, COLOR_FOREST)
    # Rising arrow in orange
    draw_connector(slide, cx - 0.08, cy + 0.05, cx + 0.08, cy - 0.07, COLOR_ORANGE, 1.2)
    # Arrow tip
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx + 0.04), Inches(cy - 0.10), Inches(0.06), Inches(0.06)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_ORANGE
    arrow.line.fill.background()

def draw_icon_lock(slide, cx, cy):
    shackle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - 0.035), Inches(cy - 0.06), Inches(0.07), Inches(0.07)
    )
    shackle.fill.background()
    shackle.line.color.rgb = COLOR_FOREST
    shackle.line.width = Pt(1)
    
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - 0.05), Inches(cy - 0.01), Inches(0.1), Inches(0.08)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = COLOR_FOREST
    body.line.fill.background()
    
    draw_oval(slide, cx, cy + 0.02, 0.01, COLOR_BG_LIGHT)

def draw_icon_globe(slide, cx, cy):
    # Earth circle
    draw_oval(slide, cx, cy, 0.075, COLOR_BG_LIGHT, line_color=COLOR_FOREST, line_width_pt=1).fill.background()
    # Continent squiggles
    draw_oval(slide, cx - 0.03, cy - 0.02, 0.025, COLOR_FOREST)
    draw_oval(slide, cx + 0.03, cy + 0.02, 0.025, COLOR_FOREST)

def create_slide_7():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # 1. Slide Index "08" (matching the reference image exactly)
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(1.5), Inches(0.5))
    tf_num = num_box.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = "08"
    p_num.font.name = 'Inter'
    p_num.font.size = Pt(22)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_ORANGE
    
    # Vertical divider on top left
    draw_line(slide, 1.45, 0.42, 1.46, 0.78, COLOR_BORDER_THIN)
    
    # Top left subheader
    sh_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.4), Inches(3.0), Inches(0.5))
    tf_sh = sh_box.text_frame
    tf_sh.word_wrap = True
    p_sh = tf_sh.paragraphs[0]
    p_sh.text = "CONNECTING INTELLIGENCE\nACROSS ORGANIZATIONS"
    p_sh.font.name = 'Inter'
    p_sh.font.size = Pt(8.5)
    p_sh.font.bold = True
    p_sh.font.color.rgb = COLOR_MUTED
    p_sh.line_spacing = 1.15
    
    # 2. Right Top Section Header: EXPECTED IMPACT
    draw_line(slide, 11.0, 0.42, 11.015, 0.78, COLOR_ORANGE)
    sol_hdr_box = slide.shapes.add_textbox(Inches(11.15), Inches(0.42), Inches(2.0), Inches(0.4))
    p_shdr = sol_hdr_box.text_frame.paragraphs[0]
    p_shdr.text = "EXPECTED IMPACT"
    p_shdr.font.name = 'Inter'
    p_shdr.font.size = Pt(10)
    p_shdr.font.bold = True
    p_shdr.font.color.rgb = COLOR_CHARCOAL
    
    # 3. Main Center Heading
    title_box = slide.shapes.add_textbox(Inches(3.0), Inches(1.0), Inches(7.333), Inches(1.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_h = tf_title.paragraphs[0]
    p_h.text = "EXPECTED IMPACT"
    p_h.font.name = 'Sora'
    p_h.font.size = Pt(44)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_FOREST
    p_h.alignment = PP_ALIGN.CENTER
    
    # Horizontal line under title
    draw_line(slide, 6.1, 1.85, 7.2, 1.87, COLOR_FOREST)
    
    # Subtitle Text
    sub_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.05), Inches(8.333), Inches(0.7))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    tf_sub.margin_left = tf_sub.margin_right = tf_sub.margin_top = tf_sub.margin_bottom = 0
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Building a connected intelligence infrastructure will create value\nfor organizations, industries, and society at large."
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = COLOR_CHARCOAL
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.line_spacing = 1.25
    
    # 4. 6-Column Grid Layout
    cols_data = [
        # (col_idx, title, desc, icon_type)
        (0, "SECURE\nCOLLABORATION", "Organizations can share and collaborate without compromising privacy or control.", 'shield_check'),
        (1, "FASTER\nINNOVATION", "Combine diverse intelligence and expertise to solve complex problems faster.", 'lightbulb'),
        (2, "STRONGER\nECOSYSTEM", "Build a trusted network of organizations working together toward common goals.", 'group'),
        (3, "GREATER\nEFFICIENCY", "Reduce duplication of efforts and optimize the use of resources.", 'chart'),
        (4, "PRIVACY &\nTRUST", "Ensure transparency, accountability, and data protection at every step.", 'lock'),
        (5, "SOCIETAL\nIMPACT", "Drive inclusive growth and create positive impact across industries and communities.", 'globe')
    ]
    
    for idx, title, desc, icon_type in cols_data:
        cx = 0.8 + idx * 1.95 + 0.85
        cy = 3.65
        
        # Circle outline with thin green border
        draw_oval(slide, cx, cy, 0.38, COLOR_BG_LIGHT, line_color=COLOR_FOREST, line_width_pt=0.8).fill.background()
        
        # Orbital accent rings around circles
        draw_oval(slide, cx, cy, 0.44, COLOR_BG_LIGHT, line_color=COLOR_BORDER_THIN, line_width_pt=0.5).fill.background()
        draw_oval(slide, cx + 0.36, cy - 0.25, 0.02, COLOR_FOREST)
        draw_oval(slide, cx - 0.36, cy + 0.25, 0.02, COLOR_FOREST)
        
        # Render internal icon
        if icon_type == 'shield_check':
            draw_icon_shield_check(slide, cx, cy)
        elif icon_type == 'lightbulb':
            draw_icon_lightbulb(slide, cx, cy)
        elif icon_type == 'group':
            draw_icon_group(slide, cx, cy)
        elif icon_type == 'chart':
            draw_icon_chart(slide, cx, cy)
        elif icon_type == 'lock':
            draw_icon_lock(slide, cx, cy)
        elif icon_type == 'globe':
            draw_icon_globe(slide, cx, cy)
            
        # Tiny line divider
        draw_line(slide, cx - 0.2, 4.25, cx + 0.2, 4.265, COLOR_FOREST)
        
        # Column title text box
        title_box = slide.shapes.add_textbox(Inches(cx - 0.95), Inches(4.5), Inches(1.9), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Inter'
        p_t.font.size = Pt(10.5)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_FOREST
        p_t.alignment = PP_ALIGN.CENTER
        p_t.line_spacing = 1.15
        
        # Column description text box
        desc_box = slide.shapes.add_textbox(Inches(cx - 0.95), Inches(5.3), Inches(1.9), Inches(0.9))
        tf_d = desc_box.text_frame
        tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_right = tf_d.margin_top = tf_d.margin_bottom = 0
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = COLOR_CHARCOAL
        p_d.alignment = PP_ALIGN.CENTER
        p_d.line_spacing = 1.25

    # 5. Bottom Quote Highlight Box (Warm Sand)
    quote_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.9)
    )
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = COLOR_SAND
    quote_card.line.fill.background()
    
    # Left circle target icon (Forest Green with Orange target inside)
    draw_oval(slide, 1.3, 6.75, 0.25, COLOR_FOREST)
    draw_oval(slide, 1.3, 6.75, 0.15, COLOR_FOREST, line_color=COLOR_ORANGE, line_width_pt=1.2)
    draw_oval(slide, 1.3, 6.75, 0.05, COLOR_ORANGE)
    
    # Left vertical orange indicator line
    draw_line(slide, 1.8, 6.45, 1.815, 7.05, COLOR_ORANGE)
    
    # Quote text box (Baskerville serif style)
    quote_box = slide.shapes.add_textbox(Inches(2.0), Inches(6.45), Inches(10.0), Inches(0.6))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    tf_q.margin_left = tf_q.margin_right = tf_q.margin_top = tf_q.margin_bottom = 0
    p_q = tf_q.paragraphs[0]
    p_q.text = ""
    
    r1 = p_q.add_run()
    r1.text = "A connected intelligence ecosystem leads to a smarter, more collaborative, and "
    r1.font.name = 'Libre Baskerville'
    r1.font.size = Pt(13)
    r1.font.color.rgb = COLOR_CHARCOAL
    
    r2 = p_q.add_run()
    r2.text = "sustainable future"
    r2.font.name = 'Libre Baskerville'
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_FOREST
    
    r3 = p_q.add_run()
    r3.text = " for everyone."
    r3.font.name = 'Libre Baskerville'
    r3.font.size = Pt(13)
    r3.font.color.rgb = COLOR_CHARCOAL

    # Save
    output_filename = "presentation_slide7_final.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    try:
        prs.save(output_path)
        print(f"Slide 7 presentation successfully generated at: {output_path}")
    except PermissionError:
        fallback_path = os.path.join(output_dir, "presentation_slide7_fixed.pptx")
        prs.save(fallback_path)
        print(f"[!] WARNING: '{output_path}' was locked/open. Saved fallback to: {fallback_path}")

if __name__ == "__main__":
    create_slide_7()
