import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ==========================================================================
# Core Configuration & Helper Functions
# ==========================================================================

# Colors
COLOR_BG = RGBColor(17, 19, 21)          # Charcoal (#111315)
COLOR_CARD = RGBColor(25, 27, 29)        # Lighter Charcoal for cards
COLOR_COPPER = RGBColor(184, 115, 51)    # Copper (#B87333)
COLOR_EMERALD = RGBColor(16, 185, 129)   # Emerald (#10B981)
COLOR_WHITE = RGBColor(245, 245, 247)    # Off-white (#F5F5F7)
COLOR_GRAY = RGBColor(134, 134, 139)     # Muted gray (#86868B)
COLOR_DARK_GRAY = RGBColor(81, 81, 84)   # Dark gray (#515154)
COLOR_BORDER = RGBColor(44, 45, 48)      # Card border (#2C2D30)

def set_slide_bg(slide):
    """Fills the slide background with solid Matte Charcoal."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, slide_num, title_text):
    """Adds a consistent premium title header with number indicator."""
    # Number box (e.g. "02")
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(1.2), Inches(0.8))
    tf_num = num_box.text_frame
    tf_num.word_wrap = True
    tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
    p_num = tf_num.paragraphs[0]
    p_num.text = slide_num
    p_num.font.name = 'Sora'
    p_num.font.size = Pt(22)
    p_num.font.bold = True
    p_num.font.color.rgb = COLOR_COPPER
    
    # Title box
    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.4), Inches(10.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Sora'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

def draw_oval(slide, cx, cy, r, color, line_color=None, line_width_pt=1):
    """Draws a circle shape given center coordinates and radius in inches."""
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, 
        Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = color
    if line_color:
        oval.line.color.rgb = line_color
        oval.line.width = Pt(line_width_pt)
    else:
        oval.line.fill.background()
    return oval

def draw_line(slide, x1, y1, x2, y2, color, width_pt=1.0):
    """Draws a straight connector line between two points in inches."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, 
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn

# ==========================================================================
# Main PPTX Builder
# ==========================================================================

def create_presentation():
    prs = Presentation()
    
    # Set to standard 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------------------------
    slide_1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_1)
    
    # Left Content Textbox (Flowing titles for no overlaps)
    left_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(7.5), Inches(4.0))
    tf_1 = left_box.text_frame
    tf_1.word_wrap = True
    tf_1.margin_left = tf_1.margin_right = tf_1.margin_top = tf_1.margin_bottom = 0
    
    # B.Tech tag
    p_tag = tf_1.paragraphs[0]
    p_tag.text = "B.TECH FINAL YEAR PROJECT PRESENTATION"
    p_tag.font.name = 'Sora'
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_COPPER
    p_tag.space_after = Pt(24)
    
    # Main Title
    p_title = tf_1.add_paragraph()
    p_title.text = "Connecting Intelligence\nAcross Organizations"
    p_title.font.name = 'Sora'
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.space_after = Pt(20)
    
    # Subtitle
    p_sub = tf_1.add_paragraph()
    p_sub.text = "A Distributed Intelligence Infrastructure for Secure AI Collaboration"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_GRAY
    
    # Thin divider line
    div_rect = slide_1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0.8), Inches(4.7), Inches(11.733), Inches(0.01)
    )
    div_rect.fill.solid()
    div_rect.fill.fore_color.rgb = COLOR_COPPER
    div_rect.line.fill.background()
    
    # Presenters block (Bottom Left)
    pres_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(3.8), Inches(2.0))
    tf_pres = pres_box.text_frame
    tf_pres.word_wrap = True
    tf_pres.margin_left = tf_pres.margin_right = tf_pres.margin_top = tf_pres.margin_bottom = 0
    
    p_plbl = tf_pres.paragraphs[0]
    p_plbl.text = "PRESENTED BY"
    p_plbl.font.name = 'Sora'
    p_plbl.font.size = Pt(10)
    p_plbl.font.bold = True
    p_plbl.font.color.rgb = COLOR_DARK_GRAY
    p_plbl.space_after = Pt(8)
    
    p_pnames = tf_pres.add_paragraph()
    p_pnames.text = "Gokul S\nGautham Krishna R Nair\nSreenidhi V\nT Maheswaran"
    p_pnames.font.name = 'Inter'
    p_pnames.font.size = Pt(13)
    p_pnames.font.color.rgb = COLOR_WHITE
    
    # Guide block (Bottom Right-Center)
    guide_box = slide_1.shapes.add_textbox(Inches(5.0), Inches(5.1), Inches(4.5), Inches(2.0))
    tf_guide = guide_box.text_frame
    tf_guide.word_wrap = True
    tf_guide.margin_left = tf_guide.margin_right = tf_guide.margin_top = tf_guide.margin_bottom = 0
    
    p_glbl = tf_guide.paragraphs[0]
    p_glbl.text = "PROJECT GUIDE"
    p_glbl.font.name = 'Sora'
    p_glbl.font.size = Pt(10)
    p_glbl.font.bold = True
    p_glbl.font.color.rgb = COLOR_DARK_GRAY
    p_glbl.space_after = Pt(8)
    
    p_gname = tf_guide.add_paragraph()
    p_gname.text = "Assistant Professor Nisha Soman"
    p_gname.font.name = 'Inter'
    p_gname.font.size = Pt(13)
    p_gname.font.bold = True
    p_gname.font.color.rgb = COLOR_WHITE
    p_gname.space_after = Pt(4)
    
    p_gdept = tf_guide.add_paragraph()
    p_gdept.text = "Department of AIML\nMarian Engineering College, Thiruvananthapuram"
    p_gdept.font.name = 'Inter'
    p_gdept.font.size = Pt(11)
    p_gdept.font.color.rgb = COLOR_GRAY
    
    # Elegant abstract visual illustration (Right side)
    # Circle positions (cx, cy, r)
    nodes_s1 = [
        (9.5, 2.2, 0.08), (11.8, 1.8, 0.06), (10.8, 3.2, 0.12),
        (12.2, 3.4, 0.07), (9.8, 4.3, 0.09), (11.5, 4.5, 0.10)
    ]
    # Connections
    draw_line(slide_1, 9.5, 2.2, 10.8, 3.2, COLOR_COPPER, 0.75)
    draw_line(slide_1, 11.8, 1.8, 10.8, 3.2, COLOR_COPPER, 0.75)
    draw_line(slide_1, 10.8, 3.2, 12.2, 3.4, COLOR_COPPER, 0.75)
    draw_line(slide_1, 10.8, 3.2, 9.8, 4.3, COLOR_COPPER, 0.75)
    draw_line(slide_1, 9.8, 4.3, 11.5, 4.5, COLOR_COPPER, 0.75)
    draw_line(slide_1, 12.2, 3.4, 11.5, 4.5, COLOR_COPPER, 0.75)
    
    # Render nodes
    for cx, cy, r in nodes_s1:
        draw_oval(slide_1, cx, cy, r, COLOR_COPPER)

    # ----------------------------------------------------------------------
    # SLIDE 2: Abstract
    # ----------------------------------------------------------------------
    slide_2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_2)
    add_header(slide_2, "02", "Abstract")
    
    # Left column text box
    abstract_box = slide_2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf_abs = abstract_box.text_frame
    tf_abs.word_wrap = True
    tf_abs.margin_left = tf_abs.margin_right = tf_abs.margin_top = tf_abs.margin_bottom = 0
    
    # Paragraph 1 (Lead paragraph with copper side border style or spacing)
    p_abs1 = tf_abs.paragraphs[0]
    p_abs1.text = "Artificial Intelligence is transforming industries through intelligent decision-making and automation. However, AI systems developed by different organizations often work independently because of privacy, security, and ownership concerns."
    p_abs1.font.name = 'Inter'
    p_abs1.font.size = Pt(17)
    p_abs1.font.color.rgb = COLOR_WHITE
    p_abs1.space_after = Pt(28)
    
    # Paragraph 2
    p_abs2 = tf_abs.add_paragraph()
    p_abs2.text = "This project proposes a Distributed Intelligence Infrastructure that enables AI systems to collaborate securely without compromising organizational control. The objective is to improve collaboration, accelerate innovation, and establish a connected AI ecosystem."
    p_abs2.font.name = 'Inter'
    p_abs2.font.size = Pt(14)
    p_abs2.font.color.rgb = COLOR_GRAY
    
    # Right column - Central collaboration illustration
    # Card background
    card_abs = slide_2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.3)
    )
    card_abs.fill.solid()
    card_abs.fill.fore_color.rgb = COLOR_CARD
    card_abs.line.color.rgb = COLOR_BORDER
    card_abs.line.width = Pt(1)
    
    # Core emerald central node
    ccx, ccy = 10.25, 3.95
    
    # Faint outer copper nodes and connections
    outer_nodes_s2 = [(8.8, 2.7), (11.7, 2.6), (9.0, 5.0), (11.5, 4.9)]
    for ox, oy in outer_nodes_s2:
        draw_line(slide_2, ox, oy, ccx, ccy, COLOR_COPPER, 0.75)
        draw_oval(slide_2, ox, oy, 0.08, COLOR_COPPER)
        
    draw_oval(slide_2, ccx, ccy, 0.20, COLOR_EMERALD)
    
    # Label on card
    lbl_box = slide_2.shapes.add_textbox(Inches(8.2), Inches(5.4), Inches(4.1), Inches(0.5))
    tf_lbl = lbl_box.text_frame
    tf_lbl.word_wrap = True
    p_lbl = tf_lbl.paragraphs[0]
    p_lbl.text = "Central Collaboration Layer"
    p_lbl.font.name = 'Sora'
    p_lbl.font.size = Pt(11)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = COLOR_GRAY
    p_lbl.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # ----------------------------------------------------------------------
    slide_3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_3)
    add_header(slide_3, "03", "Problem Statement")
    
    # Left column: Challenges text box
    prob_box = slide_3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf_prob = prob_box.text_frame
    tf_prob.word_wrap = True
    tf_prob.margin_left = tf_prob.margin_right = tf_prob.margin_top = tf_prob.margin_bottom = 0
    
    p_ph = tf_prob.paragraphs[0]
    p_ph.text = "Current Challenges"
    p_ph.font.name = 'Sora'
    p_ph.font.size = Pt(20)
    p_ph.font.bold = True
    p_ph.font.color.rgb = COLOR_WHITE
    p_ph.space_after = Pt(24)
    
    challenges = [
        "AI systems operate independently.",
        "Organizations cannot easily collaborate.",
        "Privacy prevents direct data sharing.",
        "Similar AI solutions are repeatedly developed.",
        "No common infrastructure exists for collaboration."
    ]
    for ch in challenges:
        p_ch = tf_prob.add_paragraph()
        p_ch.text = "•   " + ch
        p_ch.font.name = 'Inter'
        p_ch.font.size = Pt(14)
        p_ch.font.color.rgb = COLOR_GRAY
        p_ch.space_after = Pt(12)
        
    # Right column: Isolated silos card
    card_prob = slide_3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.3)
    )
    card_prob.fill.solid()
    card_prob.fill.fore_color.rgb = COLOR_CARD
    card_prob.line.color.rgb = COLOR_BORDER
    card_prob.line.width = Pt(1)
    
    # Isolated nodes layout (No connections)
    # Hospital, Research Lab, University, Industry, Government
    silos_coords = [
        (8.9, 2.5, "Hospital"),
        (11.6, 2.5, "Research Lab"),
        (10.25, 3.7, "University"),
        (9.1, 4.9, "Industry"),
        (11.4, 4.9, "Government")
    ]
    for sx, sy, label in silos_coords:
        # Draw red border circle representing isolated nodes
        draw_oval(slide_3, sx, sy, 0.08, COLOR_GRAY, line_color=RGBColor(239, 68, 68), line_width_pt=1.5)
        
        # Add labels centered above/below
        lbl_box = slide_3.shapes.add_textbox(Inches(sx - 1.0), Inches(sy + 0.15), Inches(2.0), Inches(0.4))
        tf_s = lbl_box.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.text = label
        p_s.font.name = 'Inter'
        p_s.font.size = Pt(10)
        p_s.font.bold = True
        p_s.font.color.rgb = COLOR_GRAY
        p_s.alignment = PP_ALIGN.CENTER
        
    # SILOS overlay text
    status_box = slide_3.shapes.add_textbox(Inches(8.2), Inches(1.9), Inches(1.5), Inches(0.4))
    tf_status = status_box.text_frame
    p_status = tf_status.paragraphs[0]
    p_status.text = "ISOLATED SILOS"
    p_status.font.name = 'Sora'
    p_status.font.size = Pt(9)
    p_status.font.bold = True
    p_status.font.color.rgb = RGBColor(239, 68, 68)
    
    # Bottom Quote
    quote_box = slide_3.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.5))
    tf_q = quote_box.text_frame
    tf_q.word_wrap = True
    p_q = tf_q.paragraphs[0]
    p_q.text = '"The next challenge in AI is enabling independent AI systems to collaborate securely."'
    p_q.font.name = 'Inter'
    p_q.font.size = Pt(14)
    p_q.font.italic = True
    p_q.font.color.rgb = COLOR_GRAY
    p_q.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------------------------
    # SLIDE 4: Proposed Solution
    # ----------------------------------------------------------------------
    slide_4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_4)
    add_header(slide_4, "04", "Proposed Solution")
    
    # Left column: Content text box
    sol_box = slide_4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.5))
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True
    tf_sol.margin_left = tf_sol.margin_right = tf_sol.margin_top = tf_sol.margin_bottom = 0
    
    p_subsol = tf_sol.paragraphs[0]
    p_subsol.text = "THE NEW PARADIGM"
    p_subsol.font.name = 'Sora'
    p_subsol.font.size = Pt(10)
    p_subsol.font.bold = True
    p_subsol.font.color.rgb = COLOR_COPPER
    p_subsol.space_after = Pt(8)
    
    p_solh = tf_sol.add_paragraph()
    p_solh.text = "Distributed Intelligence Infrastructure"
    p_solh.font.name = 'Sora'
    p_solh.font.size = Pt(22)
    p_solh.font.bold = True
    p_solh.font.color.rgb = COLOR_WHITE
    p_solh.space_after = Pt(20)
    
    sol_points = [
        "Communicate securely",
        "Coordinate tasks",
        "Share approved capabilities",
        "Preserve ownership",
        "Maintain privacy"
    ]
    for pt in sol_points:
        p_pt = tf_sol.add_paragraph()
        p_pt.text = "•   " + pt
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(14)
        p_pt.font.color.rgb = COLOR_GRAY
        p_pt.space_after = Pt(10)
        
    # Right column: Connected infrastructure card
    card_sol = slide_4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.3)
    )
    card_sol.fill.solid()
    card_sol.fill.fore_color.rgb = COLOR_CARD
    card_sol.line.color.rgb = COLOR_COPPER
    card_sol.line.width = Pt(1)
    
    # Hub center
    hcx, hcy = 10.25, 3.95
    
    # Nodes connect to central hub
    for sx, sy, label in silos_coords:
        draw_line(slide_4, sx, sy, hcx, hcy, COLOR_COPPER, 1.25)
        draw_oval(slide_4, sx, sy, 0.08, COLOR_COPPER)
        
    # Draw central DII core hub
    draw_oval(slide_4, hcx, hcy, 0.16, COLOR_EMERALD)
    
    # Label above central hub
    hub_lbl = slide_4.shapes.add_textbox(Inches(hcx - 1.2), Inches(hcy - 0.55), Inches(2.4), Inches(0.4))
    tf_hl = hub_lbl.text_frame
    p_hl = tf_hl.paragraphs[0]
    p_hl.text = "DII Core"
    p_hl.font.name = 'Sora'
    p_hl.font.size = Pt(10)
    p_hl.font.bold = True
    p_hl.font.color.rgb = COLOR_WHITE
    p_hl.alignment = PP_ALIGN.CENTER
    
    # Bottom statement
    sol_stmt = slide_4.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.733), Inches(0.5))
    tf_stmt = sol_stmt.text_frame
    p_stmt = tf_stmt.paragraphs[0]
    p_stmt.text = "Connecting intelligence without compromising trust."
    p_stmt.font.name = 'Sora'
    p_stmt.font.size = Pt(13)
    p_stmt.font.bold = True
    p_stmt.font.color.rgb = COLOR_COPPER
    p_stmt.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------------------------
    # SLIDE 5: Implementation Plan
    # ----------------------------------------------------------------------
    slide_5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_5)
    add_header(slide_5, "05", "Implementation Plan")
    
    # Horizontal timeline line
    timeline_line = slide_5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(1.0), Inches(3.6), Inches(11.333), Inches(0.02)
    )
    timeline_line.fill.solid()
    timeline_line.fill.fore_color.rgb = COLOR_BORDER
    timeline_line.line.fill.background()
    
    # 4 Stages
    stages_data = [
        ("STAGE 01", "Research", "Study AI collaboration approaches.", 2.2),
        ("STAGE 02", "Design", "Design secure communication architecture.", 5.0),
        ("STAGE 03", "Development", "Develop a working prototype.", 7.8),
        ("STAGE 04", "Evaluation", "Test collaboration, security and scalability.", 10.6)
    ]
    
    for meta, title, desc, cx in stages_data:
        # Draw timeline dot
        draw_oval(slide_5, cx, 3.6, 0.15, COLOR_COPPER)
        
        # Text box above dot (meta, title)
        stage_box = slide_5.shapes.add_textbox(Inches(cx - 1.2), Inches(2.0), Inches(2.4), Inches(1.3))
        tf_st = stage_box.text_frame
        tf_st.word_wrap = True
        
        p_m = tf_st.paragraphs[0]
        p_m.text = meta
        p_m.font.name = 'Sora'
        p_m.font.size = Pt(10)
        p_m.font.bold = True
        p_m.font.color.rgb = COLOR_COPPER
        p_m.alignment = PP_ALIGN.CENTER
        p_m.space_after = Pt(4)
        
        p_t = tf_st.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Sora'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.alignment = PP_ALIGN.CENTER
        
        # Text box below dot (description)
        desc_box = slide_5.shapes.add_textbox(Inches(cx - 1.2), Inches(4.0), Inches(2.4), Inches(1.8))
        tf_ds = desc_box.text_frame
        tf_ds.word_wrap = True
        
        p_d = tf_ds.paragraphs[0]
        p_d.text = desc
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_GRAY
        p_d.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------------------------
    # SLIDE 6: Target Users & Expected Impact
    # ----------------------------------------------------------------------
    slide_6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_6)
    add_header(slide_6, "06", "Target Users & Expected Impact")
    
    # Left Column: Target Users
    user_title_box = slide_6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5))
    p_ut = user_title_box.text_frame.paragraphs[0]
    p_ut.text = "Target Users"
    p_ut.font.name = 'Sora'
    p_ut.font.size = Pt(20)
    p_ut.font.bold = True
    p_ut.font.color.rgb = COLOR_COPPER
    
    users = [
        "Educational Institutions",
        "Research Organizations",
        "Healthcare",
        "Government",
        "Enterprises",
        "AI Developers"
    ]
    
    user_list_box = slide_6.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
    tf_u = user_list_box.text_frame
    tf_u.word_wrap = True
    tf_u.margin_left = tf_u.margin_right = tf_u.margin_top = tf_u.margin_bottom = 0
    for u in users:
        p_u = tf_u.add_paragraph() if tf_u.paragraphs[0].text else tf_u.paragraphs[0]
        p_u.text = "•   " + u
        p_u.font.name = 'Inter'
        p_u.font.size = Pt(14)
        p_u.font.color.rgb = COLOR_WHITE
        p_u.space_after = Pt(14)
        
    # Right Column: Expected Impact
    impact_title_box = slide_6.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5))
    p_it = impact_title_box.text_frame.paragraphs[0]
    p_it.text = "Expected Impact"
    p_it.font.name = 'Sora'
    p_it.font.size = Pt(20)
    p_it.font.bold = True
    p_it.font.color.rgb = COLOR_EMERALD
    
    impacts = [
        "Secure collaboration",
        "Faster innovation",
        "Knowledge sharing",
        "Reduced duplication",
        "Better utilization of AI resources"
    ]
    
    impact_list_box = slide_6.shapes.add_textbox(Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5))
    tf_i = impact_list_box.text_frame
    tf_i.word_wrap = True
    tf_i.margin_left = tf_i.margin_right = tf_i.margin_top = tf_i.margin_bottom = 0
    for im in impacts:
        p_i = tf_i.add_paragraph() if tf_i.paragraphs[0].text else tf_i.paragraphs[0]
        p_i.text = "•   " + im
        p_i.font.name = 'Inter'
        p_i.font.size = Pt(14)
        p_i.font.color.rgb = COLOR_WHITE
        p_i.space_after = Pt(14)

    # ----------------------------------------------------------------------
    # SLIDE 7: Conclusion
    # ----------------------------------------------------------------------
    slide_7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide_7)
    add_header(slide_7, "07", "Conclusion")
    
    # Left column text box
    conc_box = slide_7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
    tf_conc = conc_box.text_frame
    tf_conc.word_wrap = True
    tf_conc.margin_left = tf_conc.margin_right = tf_conc.margin_top = tf_conc.margin_bottom = 0
    
    p_c1 = tf_conc.paragraphs[0]
    p_c1.text = "Artificial Intelligence is becoming an essential technology across industries, but collaboration between independent AI systems remains a major challenge."
    p_c1.font.name = 'Inter'
    p_c1.font.size = Pt(16)
    p_c1.font.color.rgb = COLOR_WHITE
    p_c1.space_after = Pt(24)
    
    p_c2 = tf_conc.add_paragraph()
    p_c2.text = "This project proposes a Distributed Intelligence Infrastructure that enables secure collaboration while preserving privacy, ownership, and trust."
    p_c2.font.name = 'Inter'
    p_c2.font.size = Pt(14)
    p_c2.font.color.rgb = COLOR_GRAY
    p_c2.space_after = Pt(20)
    
    p_c3 = tf_conc.add_paragraph()
    p_c3.text = "By connecting intelligence across organizations, the proposed approach aims to support the next generation of collaborative AI."
    p_c3.font.name = 'Inter'
    p_c3.font.size = Pt(14)
    p_c3.font.color.rgb = COLOR_GRAY
    
    # Right column quote block card
    card_conc = slide_7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(7.3), Inches(1.8), Inches(5.2), Inches(4.3)
    )
    card_conc.fill.solid()
    card_conc.fill.fore_color.rgb = COLOR_CARD
    card_conc.line.color.rgb = COLOR_BORDER
    card_conc.line.width = Pt(1)
    
    # Quote copper bar shape
    bar_conc = slide_7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(7.3), Inches(1.8), Inches(0.08), Inches(4.3)
    )
    bar_conc.fill.solid()
    bar_conc.fill.fore_color.rgb = COLOR_COPPER
    bar_conc.line.fill.background()
    
    # Text inside quote block
    quote_text_box = slide_7.shapes.add_textbox(Inches(7.6), Inches(2.2), Inches(4.6), Inches(3.5))
    tf_qconc = quote_text_box.text_frame
    tf_qconc.word_wrap = True
    
    p_qc = tf_qconc.paragraphs[0]
    p_qc.text = '"The future of Artificial Intelligence is not only about building intelligent systems, but enabling them to work together."'
    p_qc.font.name = 'Sora'
    p_qc.font.size = Pt(20)
    p_qc.font.bold = True
    p_qc.font.italic = True
    p_qc.font.color.rgb = COLOR_WHITE
    p_qc.space_after = Pt(15)
    
    p_qc_sub = tf_qconc.add_paragraph()
    p_qc_sub.text = "— Closing Statement"
    p_qc_sub.font.name = 'Inter'
    p_qc_sub.font.size = Pt(12)
    p_qc_sub.font.color.rgb = COLOR_COPPER

    # Save presentation
    output_filename = "presentation.pptx"
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, output_filename)
    
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
